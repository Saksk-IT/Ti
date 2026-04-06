from __future__ import annotations

import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOC_ROOT = ROOT / 'output/doc/1软件应用与开发_1Web应用与开发_Sak-AI答题助手'
OUTPUT_DIR = DOC_ROOT / '01作品与答辩材料'
ASSET_DIR = DOC_ROOT / '02素材与源码'
SCREENSHOT_DIR = ASSET_DIR / '网站截图'
ILLUSTRATION_DIR = ASSET_DIR / '插图'
OUTPUT_PPTX = OUTPUT_DIR / 'Sak-AI答题助手-答辩演示PPT.pptx'

FONT_FAMILY = 'Microsoft YaHei'
EN_FONT_FAMILY = 'Aptos'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = '0F172A'
NAVY_2 = '172554'
INK = '102A43'
TEXT = '17324D'
MUTED = '62748A'
SUBTLE = '94A3B8'
ACCENT = '4F6EF7'
ACCENT_2 = '5BC0BE'
ACCENT_3 = 'D6B36A'
SURFACE = 'F7F9FC'
CARD = 'FFFFFF'
CARD_SOFT = 'F3F7FF'
BORDER = 'D9E2EF'
DARK_BORDER = '31415F'
SUCCESS = '34A58D'


def rgb(value: str) -> RGBColor:
    value = value.replace('#', '')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def emu_to_px(length, scale: int = 150) -> int:
    inches = length / 914400
    return max(1, int(round(inches * scale)))


def add_full_bleed(slide, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def make_dark_background(path: Path, seed: int = 0):
    img = Image.new('RGBA', (1600, 900), f'#{NAVY}')
    base = ImageDraw.Draw(img)
    for y in range(900):
        ratio = y / 899
        r = int(15 + (24 - 15) * ratio)
        g = int(23 + (43 - 23) * ratio)
        b = int(42 + (78 - 42) * ratio)
        base.line((0, y, 1600, y), fill=(r, g, b, 255))

    glow = Image.new('RGBA', img.size, (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    circles = [
        ((1180, 180), 320, (*rgb_tuple(ACCENT), 110)),
        ((1320, 620), 250, (*rgb_tuple(ACCENT_2), 80)),
        ((310, 150), 220, (*rgb_tuple(ACCENT_3), 50)),
        ((260 + seed * 20, 720), 260, (*rgb_tuple(ACCENT), 45)),
    ]
    for center, radius, color in circles:
        x, y = center
        gd.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color)
    glow = glow.filter(ImageFilter.GaussianBlur(90))
    img = Image.alpha_composite(img, glow)

    deco = ImageDraw.Draw(img)
    for x in range(80, 1600, 80):
        alpha = 15 if x % 160 else 28
        deco.line((x, 0, x, 900), fill=(255, 255, 255, alpha), width=1)
    for y in range(80, 900, 80):
        alpha = 12 if y % 160 else 22
        deco.line((0, y, 1600, y), fill=(255, 255, 255, alpha), width=1)

    deco.rounded_rectangle((60, 60, 1540, 840), radius=40, outline=(255, 255, 255, 22), width=2)
    deco.rounded_rectangle((85, 85, 1515, 815), radius=34, outline=(255, 255, 255, 14), width=1)
    img.save(path)



def make_light_background(path: Path, seed: int = 0):
    img = Image.new('RGBA', (1600, 900), f'#{SURFACE}')
    draw = ImageDraw.Draw(img)
    for y in range(900):
        ratio = y / 899
        start = (247, 249, 252)
        end = (241, 245, 255)
        color = tuple(int(start[i] + (end[i] - start[i]) * ratio) for i in range(3))
        draw.line((0, y, 1600, y), fill=(*color, 255))

    glow = Image.new('RGBA', img.size, (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    accent_a = (*rgb_tuple(ACCENT), 55)
    accent_b = (*rgb_tuple(ACCENT_2), 35)
    accent_c = (*rgb_tuple(ACCENT_3), 25)
    gd.ellipse((-120, -80, 500, 520), fill=accent_a)
    gd.ellipse((1040, -120, 1680, 420), fill=accent_b)
    gd.ellipse((980, 560, 1540, 980), fill=accent_c)
    glow = glow.filter(ImageFilter.GaussianBlur(70))
    img = Image.alpha_composite(img, glow)

    deco = ImageDraw.Draw(img)
    for x in range(80, 1600, 80):
        deco.line((x, 0, x, 900), fill=(71, 85, 105, 10), width=1)
    for y in range(80, 900, 80):
        deco.line((0, y, 1600, y), fill=(71, 85, 105, 8), width=1)

    for idx in range(5):
        x0 = 960 + idx * 60 + seed * 6
        deco.rounded_rectangle((x0, 110 + idx * 24, x0 + 270, 170 + idx * 24), radius=28, outline=(79, 110, 247, 28), width=2)

    img.save(path)



def rgb_tuple(value: str) -> tuple[int, int, int]:
    value = value.replace('#', '')
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)



def style_text_frame(tf, margin=0.0):
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP



def add_textbox(slide, x, y, w, h, text, *, size=18, color=TEXT, bold=False,
                font=FONT_FAMILY, align=PP_ALIGN.LEFT, italic=False, line_spacing=1.15):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    style_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return shape



def add_rich_lines(slide, x, y, w, h, lines):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    style_text_frame(tf)
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = line.get('align', PP_ALIGN.LEFT)
        p.line_spacing = line.get('line_spacing', 1.15)
        p.space_after = Pt(line.get('space_after', 4))
        p.space_before = Pt(line.get('space_before', 0))
        run = p.add_run()
        run.text = line['text']
        run.font.name = line.get('font', FONT_FAMILY)
        run.font.size = Pt(line.get('size', 18))
        run.font.bold = line.get('bold', False)
        run.font.italic = line.get('italic', False)
        run.font.color.rgb = rgb(line.get('color', TEXT))
    return shape



def add_panel(slide, x, y, w, h, *, fill=CARD, line=BORDER, radius=0.18, line_width=1.2):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape



def add_chip(slide, x, y, w, h, text, *, fill=CARD_SOFT, line='D5E0FF', text_color=ACCENT, size=12.5):
    shape = add_panel(slide, x, y, w, h, fill=fill, line=line, radius=0.35, line_width=1)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.06), w - Inches(0.24), h - Inches(0.12), text,
                size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return shape



def add_stat_card(slide, x, y, w, h, number, label, note=None, *, accent=ACCENT, dark=False):
    fill = NAVY if dark else CARD
    line = DARK_BORDER if dark else BORDER
    title_color = 'FFFFFF' if dark else accent
    label_color = 'D7DEEA' if dark else TEXT
    note_color = '9FB0C7' if dark else MUTED
    add_panel(slide, x, y, w, h, fill=fill, line=line, radius=0.22)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.5), Inches(0.48), number,
                size=24, color=title_color, bold=True, font=EN_FONT_FAMILY)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.62), w - Inches(0.5), Inches(0.34), label,
                size=12.5, color=label_color, bold=True)
    if note:
        add_textbox(slide, x + Inches(0.25), y + Inches(0.95), w - Inches(0.5), Inches(0.28), note,
                    size=10.5, color=note_color)



def section_header(slide, kicker, title, subtitle, *, dark=False):
    if dark:
        kicker_color = 'C7D2FE'
        title_color = 'FFFFFF'
        sub_color = 'D9E5F5'
        line_color = '4F6EF7'
    else:
        kicker_color = ACCENT
        title_color = INK
        sub_color = MUTED
        line_color = ACCENT
    add_textbox(slide, Inches(0.72), Inches(0.52), Inches(3.0), Inches(0.28), kicker,
                size=12, color=kicker_color, bold=True, font=EN_FONT_FAMILY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(0.88), Inches(0.76), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(line_color)
    line.line.fill.background()
    add_textbox(slide, Inches(0.72), Inches(1.0), Inches(6.8), Inches(0.58), title,
                size=28, color=title_color, bold=True)
    add_textbox(slide, Inches(0.72), Inches(1.56), Inches(7.8), Inches(0.42), subtitle,
                size=13.5, color=sub_color)



def add_footer(slide, index: int, *, dark=False):
    left_color = 'D6DFEC' if dark else SUBTLE
    right_color = 'F8FAFC' if dark else TEXT
    add_textbox(slide, Inches(0.72), Inches(7.08), Inches(3.2), Inches(0.2), 'Sak-AI答题助手｜答辩演示',
                size=9.5, color=left_color, font=EN_FONT_FAMILY)
    add_textbox(slide, Inches(12.45), Inches(7.0), Inches(0.38), Inches(0.24), f'{index:02d}',
                size=11, color=right_color, bold=True, align=PP_ALIGN.RIGHT, font=EN_FONT_FAMILY)



def prep_image(src: Path, out_dir: Path, name: str, width_px: int, height_px: int,
               *, mode='contain', crop=None, background='FFFFFF') -> Path:
    img = Image.open(src).convert('RGB')
    if crop:
        left, top, right, bottom = crop
        box = (
            int(img.width * left),
            int(img.height * top),
            int(img.width * right),
            int(img.height * bottom),
        )
        img = img.crop(box)
    if mode == 'cover':
        out = ImageOps.fit(img, (width_px, height_px), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
    else:
        out = Image.new('RGB', (width_px, height_px), f'#{background}')
        copy = img.copy()
        copy.thumbnail((width_px, height_px), Image.Resampling.LANCZOS)
        x = (width_px - copy.width) // 2
        y = (height_px - copy.height) // 2
        out.paste(copy, (x, y))
    out_path = out_dir / name
    out.save(out_path, quality=95)
    return out_path



def add_browser_frame(slide, x, y, w, h, image_path: Path, label: str, *, crop=None, mode='contain', note=None, temp_dir: Path | None = None):
    add_panel(slide, x, y, w, h, fill=CARD, line=BORDER, radius=0.18)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), Inches(0.26))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb('EEF3FB')
    bar.line.fill.background()
    try:
        bar.adjustments[0] = 0.3
    except Exception:
        pass
    dot_y = y + Inches(0.145)
    for idx, color in enumerate(('FF6B6B', 'F7B955', '34C759')):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.17 + idx * 0.12), dot_y, Inches(0.06), Inches(0.06))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(color)
        dot.line.fill.background()
    add_textbox(slide, x + Inches(0.58), y + Inches(0.1), w - Inches(0.82), Inches(0.18), label,
                size=10, color=MUTED, bold=True)

    if temp_dir is None:
        temp_dir = Path(tempfile.gettempdir())
    img_w = emu_to_px(w - Inches(0.18))
    img_h = emu_to_px(h - Inches(0.58))
    fitted = prep_image(image_path, temp_dir, f'{label}-{img_w}x{img_h}.png', img_w, img_h, mode=mode, crop=crop)
    slide.shapes.add_picture(str(fitted), x + Inches(0.09), y + Inches(0.4), width=w - Inches(0.18), height=h - Inches(0.5))
    if note:
        add_textbox(slide, x + Inches(0.16), y + h - Inches(0.35), w - Inches(0.32), Inches(0.18), note,
                    size=9.5, color=SUBTLE)



def add_phone_frame(slide, x, y, w, h, image_path: Path, *, crop=None, temp_dir: Path | None = None):
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    body.fill.solid()
    body.fill.fore_color.rgb = rgb('0B1120')
    body.line.color.rgb = rgb('283548')
    body.line.width = Pt(1.5)
    try:
        body.adjustments[0] = 0.22
    except Exception:
        pass

    screen_x = x + Inches(0.12)
    screen_y = y + Inches(0.16)
    screen_w = w - Inches(0.24)
    screen_h = h - Inches(0.32)
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, screen_x, screen_y, screen_w, screen_h)
    screen.fill.solid()
    screen.fill.fore_color.rgb = rgb('FFFFFF')
    screen.line.fill.background()
    try:
        screen.adjustments[0] = 0.18
    except Exception:
        pass
    notch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.58), y + Inches(0.08), w - Inches(1.16), Inches(0.08))
    notch.fill.solid()
    notch.fill.fore_color.rgb = rgb('172554')
    notch.line.fill.background()
    try:
        notch.adjustments[0] = 0.4
    except Exception:
        pass

    if temp_dir is None:
        temp_dir = Path(tempfile.gettempdir())
    img_w = emu_to_px(screen_w)
    img_h = emu_to_px(screen_h)
    fitted = prep_image(image_path, temp_dir, f'phone-{img_w}x{img_h}.png', img_w, img_h, mode='cover', crop=crop)
    slide.shapes.add_picture(str(fitted), screen_x, screen_y, width=screen_w, height=screen_h)



def add_feature_card(slide, x, y, w, h, title, body, index, *, accent=ACCENT):
    add_panel(slide, x, y, w, h, fill=CARD, line=BORDER, radius=0.18)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.34), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(accent)
    badge.line.fill.background()
    add_textbox(slide, x + Inches(0.18), y + Inches(0.19), Inches(0.34), Inches(0.18), str(index),
                size=12, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER, font=EN_FONT_FAMILY)
    add_textbox(slide, x + Inches(0.62), y + Inches(0.16), w - Inches(0.82), Inches(0.28), title,
                size=14, color=INK, bold=True)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.56), w - Inches(0.36), h - Inches(0.72), body,
                size=11.2, color=MUTED)



def add_flow_node(slide, x, y, w, h, title, body, step, *, accent=ACCENT):
    add_panel(slide, x, y, w, h, fill=CARD, line=BORDER, radius=0.16)
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = rgb(accent)
    strip.line.fill.background()
    add_textbox(slide, x + Inches(0.18), y + Inches(0.2), Inches(0.32), Inches(0.22), f'{step:02d}',
                size=11.5, color=accent, bold=True, font=EN_FONT_FAMILY)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.52), w - Inches(0.36), Inches(0.3), title,
                size=14, color=INK, bold=True)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.92), w - Inches(0.36), h - Inches(1.08), body,
                size=10.8, color=MUTED)



def add_chevron(slide, x, y, w, h, color='D5E3FF'):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape



def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = 'Sak-AI答题助手-答辩演示PPT'
    prs.core_properties.subject = '软件应用与开发 - Web 应用与开发'
    prs.core_properties.author = 'Codex'
    prs.core_properties.company = 'Sak-AI答题助手'

    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory(prefix='sak_ai_ppt_') as tmp:
        temp_dir = Path(tmp)
        bg_cover = temp_dir / 'bg_cover.png'
        bg_dark = temp_dir / 'bg_dark.png'
        bg_light = temp_dir / 'bg_light.png'
        bg_light_2 = temp_dir / 'bg_light_2.png'
        make_dark_background(bg_cover, seed=1)
        make_dark_background(bg_dark, seed=2)
        make_light_background(bg_light, seed=0)
        make_light_background(bg_light_2, seed=2)


        # Slide 1: Cover
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_cover)
        add_chip(slide, Inches(0.72), Inches(0.58), Inches(2.3), Inches(0.34), 'Web 应用为主体', fill='1E2D4C', line='31415F', text_color='DDE7FF')
        add_rich_lines(
            slide,
            Inches(0.72), Inches(1.18), Inches(5.5), Inches(1.4),
            [
                {'text': 'Sak-AI答题助手', 'size': 29, 'color': 'FFFFFF', 'bold': True},
                {'text': '面向备考学习与题库运营场景的单仓全栈 Web 应用', 'size': 15, 'color': 'D9E5F5', 'bold': False, 'space_after': 8},
                {'text': '题库广场｜个人题库｜刷题练习｜模拟考试｜数据中心｜论坛互动｜后台管理', 'size': 12.5, 'color': 'B8C7DC'},
            ],
        )
        add_textbox(slide, Inches(0.72), Inches(3.05), Inches(4.6), Inches(0.56),
                    '不是单一刷题页面，而是围绕“题库沉淀—训练—复盘—运营”构建的完整业务系统。',
                    size=16, color='ECF3FF', line_spacing=1.25)
        add_stat_card(slide, Inches(0.72), Inches(4.18), Inches(1.72), Inches(1.22), '12+', '业务模块', 'Blueprint 按边界拆分', dark=True)
        add_stat_card(slide, Inches(2.58), Inches(4.18), Inches(1.72), Inches(1.22), '2 端', '共享语义', 'Web + 微信原生小程序', accent=ACCENT_2, dark=True)
        add_stat_card(slide, Inches(4.44), Inches(4.18), Inches(1.72), Inches(1.22), 'Docker', '可复现部署', 'web / worker / postgres / redis', accent=ACCENT_3, dark=True)
        add_chip(slide, Inches(0.72), Inches(5.82), Inches(2.6), Inches(0.34), '赛道：软件应用与开发｜Web 应用与开发', fill='172554', line='31415F', text_color='D9E5F5', size=11)
        add_chip(slide, Inches(3.45), Inches(5.82), Inches(2.55), Inches(0.34), '正式访问：https://saksk.top', fill='172554', line='31415F', text_color='D9E5F5', size=11)

        add_browser_frame(slide, Inches(7.12), Inches(0.98), Inches(5.54), Inches(3.14), SCREENSHOT_DIR / '01-首页.png', '首页工作台', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(7.58), Inches(4.26), Inches(3.65), Inches(2.38), SCREENSHOT_DIR / '05-数据中心.png', '数据中心', mode='cover', crop=(0, 0, 1, 0.32), temp_dir=temp_dir)
        add_phone_frame(slide, Inches(11.34), Inches(3.82), Inches(1.62), Inches(2.82), SCREENSHOT_DIR / '09-首页-移动端.png', crop=(0, 0, 1, 0.28), temp_dir=temp_dir)
        add_footer(slide, 1, dark=True)

        # Slide 2: Problem insight
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light)
        section_header(slide, 'PAIN POINTS', '问题背景与需求洞察', '从单点刷题工具走向完整学习闭环，是本作品的切入点。')
        add_feature_card(slide, Inches(0.72), Inches(2.05), Inches(3.8), Inches(1.48), '资源分散', '公共题库、个人题库、课程练习常常分散存在，难以沉淀为持续可复用的学习资产。', 1)
        add_feature_card(slide, Inches(0.72), Inches(3.72), Inches(3.8), Inches(1.48), '训练断层', '用户能刷题，但难以把练习、考试、错题复盘和阶段统计串成一条连续的成长链路。', 2, accent=ACCENT_2)
        add_feature_card(slide, Inches(0.72), Inches(5.39), Inches(3.8), Inches(1.48), '管理割裂', '教师或运营侧需要统一的内容管理、数据反馈与社区支撑，传统工具往往只覆盖学习端。', 3, accent=ACCENT_3)
        add_panel(slide, Inches(4.78), Inches(2.05), Inches(7.8), Inches(4.82), fill='FFFFFF', line='D7E2F2', radius=0.18)
        add_chip(slide, Inches(5.08), Inches(2.32), Inches(1.48), Inches(0.32), 'Sak-AI 的回答', fill='EDF3FF', line='D5E0FF', text_color=ACCENT, size=11.5)
        add_rich_lines(
            slide,
            Inches(5.08), Inches(2.78), Inches(6.8), Inches(2.9),
            [
                {'text': '以 Web 应用为核心入口，承载完整业务流程。', 'size': 18, 'color': INK, 'bold': True},
                {'text': '• 统一承接题库广场、个人题库、刷题、考试、数据中心、论坛与后台。', 'size': 13, 'color': TEXT},
                {'text': '• 通过单仓全栈架构，让学习端与运营端共用一套后端与数据语义。', 'size': 13, 'color': TEXT},
                {'text': '• 以微信原生小程序补足移动高频场景，而不是另起一套割裂系统。', 'size': 13, 'color': TEXT},
                {'text': '• 以 AI 题目解析增强讲解与复盘体验，提升持续学习效率。', 'size': 13, 'color': TEXT},
            ],
        )
        add_stat_card(slide, Inches(5.08), Inches(5.72), Inches(2.05), Inches(0.92), 'Web', '主体交互端', 'SSR + 原生 JS')
        add_stat_card(slide, Inches(7.31), Inches(5.72), Inches(2.05), Inches(0.92), 'Mini', '移动延展端', '微信原生 TS', accent=ACCENT_2)
        add_stat_card(slide, Inches(9.54), Inches(5.72), Inches(2.05), Inches(0.92), 'AI', '辅助增强', '题目解析与学习说明', accent=ACCENT_3)
        add_footer(slide, 2)

        # Slide 3: positioning
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light_2)
        section_header(slide, 'POSITIONING', '作品定位：一套系统，服务学习者与管理者', 'Web 为主、小程序为辅，共享同一后端语义与学习数据。')
        add_panel(slide, Inches(0.72), Inches(2.0), Inches(3.65), Inches(4.9), fill='FFFFFF', line=BORDER)
        add_chip(slide, Inches(1.02), Inches(2.26), Inches(1.1), Inches(0.32), '学习端', fill='EEF4FF', line='D7E2F2', text_color=ACCENT)
        add_rich_lines(
            slide,
            Inches(1.02), Inches(2.7), Inches(3.0), Inches(3.8),
            [
                {'text': '题库浏览与沉淀', 'size': 17, 'bold': True, 'color': INK},
                {'text': '公共题库 + 个人题库双线并行，支持加入、整理与长期复用。', 'size': 12.6, 'color': MUTED},
                {'text': '练习与考试闭环', 'size': 17, 'bold': True, 'color': INK, 'space_before': 10},
                {'text': '错题、收藏、用户答案、考试记录与学习进度持续沉淀。', 'size': 12.6, 'color': MUTED},
                {'text': '可视化复盘', 'size': 17, 'bold': True, 'color': INK, 'space_before': 10},
                {'text': '数据中心用趋势图、热力图与能力画像，让反馈可见。', 'size': 12.6, 'color': MUTED},
            ],
        )
        add_panel(slide, Inches(4.58), Inches(2.0), Inches(4.25), Inches(4.9), fill='F8FBFF', line='D7E2F2')
        add_chip(slide, Inches(4.88), Inches(2.26), Inches(1.3), Inches(0.32), '共享语义核心', fill='EAF2FF', line='D5E0FF', text_color=ACCENT)
        flow_items = [
            ('题库资源', '公共题库 / 个人题库'),
            ('刷题训练', '答案、收藏、错题、进度'),
            ('考试评估', '模拟考试与结果沉淀'),
            ('数据复盘', '趋势、热力图、能力画像'),
        ]
        for idx, (title, desc) in enumerate(flow_items):
            cy = Inches(2.88 + idx * 0.88)
            pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.98), cy, Inches(1.15), Inches(0.42))
            pill.fill.solid(); pill.fill.fore_color.rgb = rgb('1E293B'); pill.line.fill.background()
            try: pill.adjustments[0] = 0.38
            except Exception: pass
            add_textbox(slide, Inches(5.0), cy + Inches(0.08), Inches(1.1), Inches(0.2), f'{idx+1:02d}', size=13, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER, font=EN_FONT_FAMILY)
            add_textbox(slide, Inches(6.34), cy - Inches(0.02), Inches(1.55), Inches(0.2), title, size=13.2, color=INK, bold=True)
            add_textbox(slide, Inches(6.34), cy + Inches(0.2), Inches(2.05), Inches(0.18), desc, size=10.6, color=MUTED)
            if idx < len(flow_items) - 1:
                add_chevron(slide, Inches(5.44), cy + Inches(0.46), Inches(0.22), Inches(0.28))
        add_panel(slide, Inches(9.02), Inches(2.0), Inches(3.59), Inches(4.9), fill='FFFFFF', line=BORDER)
        add_chip(slide, Inches(9.32), Inches(2.26), Inches(1.1), Inches(0.32), '运营端', fill='EEF9F6', line='CFEDE8', text_color=SUCCESS)
        add_rich_lines(
            slide,
            Inches(9.32), Inches(2.72), Inches(2.85), Inches(3.95),
            [
                {'text': '论坛互动', 'size': 17, 'bold': True, 'color': INK},
                {'text': '支持帖子浏览、交流讨论与学习社区沉淀。', 'size': 12.4, 'color': MUTED},
                {'text': '后台管理', 'size': 17, 'bold': True, 'color': INK, 'space_before': 10},
                {'text': '题目、科目、用户与系统数据可统一查看与管理。', 'size': 12.4, 'color': MUTED},
                {'text': '移动延展', 'size': 17, 'bold': True, 'color': INK, 'space_before': 10},
                {'text': '微信小程序基于同一后端语义，承接随时随地学习场景。', 'size': 12.4, 'color': MUTED},
            ],
        )
        add_footer(slide, 3)

        # Slide 4: architecture
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_dark)
        section_header(slide, 'ARCHITECTURE', '系统架构：单仓全栈、双端共享、容器化可复现', '不是前后端分离多仓拼装，而是一套围绕业务模块组织的工程系统。', dark=True)
        add_browser_frame(slide, Inches(7.28), Inches(1.96), Inches(5.4), Inches(3.22), ILLUSTRATION_DIR / '系统架构图.png', '系统架构图', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(7.28), Inches(5.28), Inches(5.4), Inches(1.42), ILLUSTRATION_DIR / '核心流程图.png', '核心用户流程', mode='contain', temp_dir=temp_dir)
        add_feature_card(slide, Inches(0.72), Inches(2.06), Inches(5.98), Inches(0.96), '应用层：Web + 微信原生小程序', 'Web 采用 Flask + Jinja 服务端渲染，小程序采用微信原生 TS，两个端共享同一后端语义。', 1, accent='7DD3FC')
        add_feature_card(slide, Inches(0.72), Inches(3.18), Inches(5.98), Inches(0.96), '服务层：Flask 应用工厂 + Blueprint 模块化', '已注册 auth、main、quiz、exam、user、chat、notifications、popups、coding、user_bank、forum、admin 等 12 个模块。', 2, accent='A7F3D0')
        add_feature_card(slide, Inches(0.72), Inches(4.30), Inches(5.98), Inches(0.96), '数据层：PostgreSQL + Redis + RQ', 'PostgreSQL 存储题库与学习数据，Redis 承担缓存/限流/队列，RQ Worker 处理后台任务。', 3, accent='FDE68A')
        add_feature_card(slide, Inches(0.72), Inches(5.42), Inches(5.98), Inches(0.96), '交付层：Docker Compose 开发编排', 'compose.dev.yml 统一拉起 web、worker、postgres、redis、backup，便于演示、复现与现场切换。', 4, accent='C4B5FD')
        add_footer(slide, 4, dark=True)

        # Slide 5: capability loop
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light)
        section_header(slide, 'CAPABILITIES', '核心功能闭环：从题库沉淀到数据复盘', '围绕学习者与运营者共同构建完整业务流程。')
        nodes = [
            ('题库广场', '浏览系统题库与公开题库，快速进入学习入口。'),
            ('个人题库', '沉淀自建/已加入题库，形成个人长期资源库。'),
            ('刷题训练', '支持题目练习、收藏、错题与答案记录。'),
            ('模拟考试', '按公共题库或个人题库发起考试训练。'),
            ('数据中心', '趋势图、热力图、能力画像帮助用户复盘。'),
            ('论坛 / 后台', '支持交流互动与系统运营管理，形成支撑面。'),
        ]
        start_x = Inches(0.72)
        y = Inches(2.35)
        node_w = Inches(1.92)
        gap = Inches(0.16)
        for idx, (title, body) in enumerate(nodes):
            x = start_x + idx * (node_w + gap)
            add_flow_node(slide, x, y, node_w, Inches(2.16), title, body, idx + 1, accent=(ACCENT if idx % 3 == 0 else ACCENT_2 if idx % 3 == 1 else ACCENT_3))
            if idx < len(nodes) - 1:
                add_chevron(slide, x + node_w + Inches(0.02), y + Inches(0.86), Inches(0.12), Inches(0.36), color='DCE7FF')
        add_panel(slide, Inches(1.18), Inches(5.25), Inches(11.0), Inches(1.08), fill='FFFFFF', line='D7E2F2', radius=0.18)
        add_textbox(slide, Inches(1.48), Inches(5.5), Inches(10.4), Inches(0.28),
                    '闭环价值：让“资源组织 → 高频训练 → 阶段评估 → 可视化复盘 → 社区互动 / 运营支撑”在同一系统中连续发生。',
                    size=15.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_footer(slide, 5)

        # Slide 6: showcase A
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light_2)
        section_header(slide, 'SHOWCASE A', '真实界面展示：学习入口与题库资产', '以下截图均来自作品真实运行界面。')
        add_browser_frame(slide, Inches(0.72), Inches(2.0), Inches(5.95), Inches(2.23), SCREENSHOT_DIR / '01-首页.png', '首页工作台', mode='cover', crop=(0.06, 0.02, 0.96, 0.82), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.95), Inches(2.0), Inches(5.66), Inches(2.23), SCREENSHOT_DIR / '02-题库广场.png', '题库广场', mode='cover', crop=(0.05, 0.02, 0.98, 0.84), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(0.72), Inches(4.45), Inches(5.95), Inches(2.23), SCREENSHOT_DIR / '03-题库名片详情.png', '题库详情（名片页）', mode='cover', crop=(0.05, 0.02, 0.98, 0.86), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.95), Inches(4.45), Inches(5.66), Inches(2.23), SCREENSHOT_DIR / '04-我的题库.png', '我的题库', mode='cover', crop=(0.05, 0.02, 0.98, 0.86), temp_dir=temp_dir)
        add_footer(slide, 6)

        # Slide 7: quiz + AI
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light)
        section_header(slide, 'QUIZ + AI', '刷题、AI 解析与 AI 判分', '重点展示答题工作台、即时解析与主观题判分配置。')
        add_browser_frame(slide, Inches(0.72), Inches(1.98), Inches(7.18), Inches(4.78), SCREENSHOT_DIR / '10-刷题AI解析.png', '刷题工作台 + AI解析', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(8.16), Inches(1.98), Inches(4.47), Inches(2.6), SCREENSHOT_DIR / '11-刷题AI判分设置.png', '主观题 AI 判分设置', mode='contain', temp_dir=temp_dir)
        add_panel(slide, Inches(8.16), Inches(4.82), Inches(4.47), Inches(1.94), fill='FFFFFF', line=BORDER)
        add_chip(slide, Inches(8.44), Inches(5.08), Inches(1.2), Inches(0.32), '模块亮点', fill='EEF4FF', line='D7E2F2', text_color=ACCENT)
        add_rich_lines(
            slide,
            Inches(8.44), Inches(5.5), Inches(3.8), Inches(1.0),
            [
                {'text': '• 支持题目列表、答题卡片与解析联动。', 'size': 12.4, 'color': TEXT},
                {'text': '• AI 解析可按题触发，辅助理解知识点。', 'size': 12.4, 'color': TEXT},
                {'text': '• 主观题支持有答即对 / AI 判分 / 自评模式。', 'size': 12.4, 'color': TEXT},
            ],
        )
        add_footer(slide, 7)

        # Slide 8: forum
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light_2)
        section_header(slide, 'FORUM', '论坛模块：社区首页到帖子详情', '不仅有帖子列表，也有完整的帖子阅读、目录、标签与互动结构。')
        add_browser_frame(slide, Inches(0.72), Inches(2.0), Inches(5.86), Inches(4.72), SCREENSHOT_DIR / '07-论坛首页.png', '论坛首页', mode='cover', crop=(0.03, 0.02, 0.98, 0.88), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.86), Inches(2.0), Inches(5.78), Inches(4.72), SCREENSHOT_DIR / '12-论坛帖子详情.png', '论坛帖子详情', mode='contain', temp_dir=temp_dir)
        add_footer(slide, 8)

        # Slide 9: chat
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_dark)
        section_header(slide, 'CHAT', '聊天模块：会话、搜索与互动内容承接', '支持从资料页或系统入口进入聊天，会话中可承接题目、图片、语音与帖子转发。', dark=True)
        add_browser_frame(slide, Inches(0.72), Inches(1.98), Inches(7.4), Inches(4.82), SCREENSHOT_DIR / '13-聊天界面.png', '聊天界面', mode='contain', temp_dir=temp_dir)
        add_panel(slide, Inches(8.38), Inches(2.08), Inches(4.2), Inches(4.62), fill='18263F', line='324362', radius=0.18)
        add_chip(slide, Inches(8.7), Inches(2.34), Inches(1.32), Inches(0.32), '交互能力', fill='243656', line='39527B', text_color='D9E5F5')
        add_rich_lines(
            slide,
            Inches(8.7), Inches(2.82), Inches(3.46), Inches(2.5),
            [
                {'text': '搜索用户后可直接发起讨论。', 'size': 13, 'color': 'EAF1FF', 'bold': True},
                {'text': '• 会话列表区分消息与互动。', 'size': 12.4, 'color': 'D9E5F5'},
                {'text': '• 消息类型覆盖图片、语音、题目与帖子转发。', 'size': 12.4, 'color': 'D9E5F5'},
                {'text': '• 与论坛、个人资料等模块形成联动。', 'size': 12.4, 'color': 'D9E5F5'},
            ],
        )
        add_stat_card(slide, Inches(8.7), Inches(5.78), Inches(1.72), Inches(0.9), 'IM', '消息流', '会话 / 搜索 / 发送', dark=True)
        add_stat_card(slide, Inches(10.56), Inches(5.78), Inches(1.72), Inches(0.9), 'Link', '内容互通', '题目 / 帖子 / 图片', accent=ACCENT_2, dark=True)
        add_footer(slide, 9, dark=True)

        # Slide 10: import workflow
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light)
        section_header(slide, 'IMPORT', '个人题目导入：题目管理 + Word 工作区', '个人题库不只是存储入口，还提供完整的题目导入与维护流程。')
        add_browser_frame(slide, Inches(0.72), Inches(2.0), Inches(5.98), Inches(4.72), SCREENSHOT_DIR / '17-个人题目管理.png', '个人题目管理', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.94), Inches(2.0), Inches(5.7), Inches(4.72), SCREENSHOT_DIR / '18-个人题目Word导入.png', 'Word 导入工作区', mode='contain', temp_dir=temp_dir)
        add_footer(slide, 10)

        # Slide 11: admin
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light_2)
        section_header(slide, 'ADMIN', '后台管理：总览、内容治理与 AI 配置', '后台不仅有仪表盘，还覆盖题目治理、论坛管理与系统级 AI 配置。')
        add_browser_frame(slide, Inches(0.72), Inches(2.0), Inches(5.95), Inches(2.23), SCREENSHOT_DIR / '08-后台仪表盘.png', '后台仪表盘', mode='cover', crop=(0.04, 0.02, 0.98, 0.85), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.95), Inches(2.0), Inches(5.66), Inches(2.23), SCREENSHOT_DIR / '14-后台题目管理.png', '后台题目管理', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(0.72), Inches(4.45), Inches(5.95), Inches(2.23), SCREENSHOT_DIR / '15-后台论坛管理.png', '后台论坛管理', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(6.95), Inches(4.45), Inches(5.66), Inches(2.23), SCREENSHOT_DIR / '16-后台AI配置.png', '后台 AI 配置', mode='contain', temp_dir=temp_dir)
        add_footer(slide, 11)

        # Slide 12: mobile synergy
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_dark)
        section_header(slide, 'CROSS-END', '双端协同：Web 主体 + 小程序延展', '展示主体是 Web；移动端延展同样复用同一套后端与数据语义。', dark=True)
        add_phone_frame(slide, Inches(0.95), Inches(1.92), Inches(2.26), Inches(4.82), SCREENSHOT_DIR / '09-首页-移动端.png', crop=(0, 0, 1, 0.33), temp_dir=temp_dir)
        add_browser_frame(slide, Inches(3.62), Inches(2.2), Inches(4.46), Inches(3.0), SCREENSHOT_DIR / '01-首页.png', 'Web 首页（响应式壳层）', mode='contain', temp_dir=temp_dir)
        add_panel(slide, Inches(8.38), Inches(2.06), Inches(4.25), Inches(4.7), fill='18263F', line='324362', radius=0.18)
        add_chip(slide, Inches(8.7), Inches(2.34), Inches(1.45), Inches(0.32), '共享关键语义', fill='243656', line='39527B', text_color='D9E5F5')
        add_rich_lines(
            slide,
            Inches(8.7), Inches(2.82), Inches(3.45), Inches(2.8),
            [
                {'text': '收藏 / 错题 / 用户答案 / 用户进度 / 模拟考试', 'size': 16, 'color': 'FFFFFF', 'bold': True},
                {'text': '• Web 端承担主要交互与展示，是参赛主体。', 'size': 12.8, 'color': 'D9E5F5'},
                {'text': '• 小程序按高频移动学习场景延展，减少语义割裂。', 'size': 12.8, 'color': 'D9E5F5'},
                {'text': '• 统一经过后端与接口封装，降低维护与演示成本。', 'size': 12.8, 'color': 'D9E5F5'},
            ],
        )
        add_stat_card(slide, Inches(8.7), Inches(5.72), Inches(1.74), Inches(0.92), 'SSR', 'Web 形态', 'Flask + Jinja', dark=True)
        add_stat_card(slide, Inches(10.55), Inches(5.72), Inches(1.74), Inches(0.92), 'JWT', 'API / 小程序', 'Bearer Token', accent=ACCENT_2, dark=True)
        add_footer(slide, 12, dark=True)

        # Slide 13: innovation
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light)
        section_header(slide, 'INNOVATION', '关键亮点：工程化表达与应用价值并重', '在题库系统常见能力之外，重点强化共享语义、可复现交付与学习辅助增强。')
        add_feature_card(slide, Inches(0.72), Inches(2.22), Inches(3.9), Inches(3.94), '共享数据语义', '公共题库、个人题库、题目、收藏、错题、答案、进度、考试等能力在 Web 与小程序两端保持同一套后端语义，降低维护复杂度，也让学习记录真正可连续。', 1)
        add_feature_card(slide, Inches(4.72), Inches(2.22), Inches(3.9), Inches(3.94), 'Docker 化可复现', '开发环境默认通过 compose.dev.yml 统一拉起 web、worker、postgres、redis、backup。无论课堂演示、现场答辩还是后续交付，都能更稳定地复现运行环境。', 2, accent=ACCENT_2)
        add_feature_card(slide, Inches(8.72), Inches(2.22), Inches(3.9), Inches(3.94), 'AI 学习辅助增强', '围绕刷题流程补充 AI 解析、主观题判分与后台 AI 配置能力，让系统更贴近教学、练习与复盘的真实场景。', 3, accent=ACCENT_3)
        add_footer(slide, 13)

        # Slide 14: validation
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_light_2)
        section_header(slide, 'VALIDATION', '工程验证与答辩可演示性', '答辩现场不仅要讲清功能，还要展示项目具备稳定运行与工程交付能力。')
        add_stat_card(slide, Inches(0.72), Inches(2.1), Inches(2.4), Inches(1.34), '8000', 'Web 暴露端口', 'compose.dev.yml：flask run --reload')
        add_stat_card(slide, Inches(3.36), Inches(2.1), Inches(2.4), Inches(1.34), '/api/ping', '健康检查', '支持基础与 deep=1 检查', accent=ACCENT_2)
        add_stat_card(slide, Inches(6.0), Inches(2.1), Inches(2.4), Inches(1.34), 'PostgreSQL 16', '开发默认数据库', 'Redis + RQ 协同运行', accent=ACCENT_3)
        add_stat_card(slide, Inches(8.64), Inches(2.1), Inches(3.0), Inches(1.34), '5 项服务', '容器编排', 'web / worker / postgres / redis / backup')
        add_panel(slide, Inches(0.72), Inches(3.82), Inches(5.75), Inches(2.9), fill='FFFFFF', line=BORDER)
        add_chip(slide, Inches(1.02), Inches(4.08), Inches(1.25), Inches(0.32), '推荐演示顺序', fill='EEF4FF', line='D7E2F2', text_color=ACCENT)
        add_rich_lines(
            slide,
            Inches(1.02), Inches(4.54), Inches(4.95), Inches(1.82),
            [
                {'text': '首页 → 题库广场 → 题库详情 → 刷题/AI → 论坛帖子 → 聊天 → 个人导入 → 后台管理', 'size': 15.2, 'color': INK, 'bold': True},
                {'text': '正式访问地址：https://saksk.top', 'size': 12.6, 'color': TEXT},
                {'text': '本地演示备用：http://127.0.0.1:8000', 'size': 12.6, 'color': TEXT},
            ],
        )
        add_panel(slide, Inches(6.72), Inches(3.82), Inches(5.9), Inches(2.9), fill='FFFFFF', line=BORDER)
        add_chip(slide, Inches(7.02), Inches(4.08), Inches(1.25), Inches(0.32), '工程事实', fill='EEF9F6', line='CFEDE8', text_color=SUCCESS)
        add_rich_lines(
            slide,
            Inches(7.02), Inches(4.5), Inches(5.0), Inches(1.9),
            [
                {'text': '• Web 是主体交互端，采用 Flask + Jinja SSR。', 'size': 13, 'color': TEXT},
                {'text': '• 小程序为微信原生工程，作为移动场景延展。', 'size': 13, 'color': TEXT},
                {'text': '• 应用工厂 + Blueprint 模块化，便于后续扩展。', 'size': 13, 'color': TEXT},
                {'text': '• 新接口遵循 status / code / data / message 信封结构。', 'size': 13, 'color': TEXT},
            ],
        )
        add_footer(slide, 14)

        # Slide 15: closing
        slide = prs.slides.add_slide(blank)
        add_full_bleed(slide, bg_cover)
        add_chip(slide, Inches(0.72), Inches(0.62), Inches(2.0), Inches(0.34), 'SUMMARY', fill='1A2845', line='31415F', text_color='DDE7FF', size=11.5)
        add_rich_lines(
            slide,
            Inches(0.72), Inches(1.22), Inches(6.1), Inches(2.0),
            [
                {'text': 'Sak-AI答题助手', 'size': 30, 'bold': True, 'color': 'FFFFFF'},
                {'text': '让题库沉淀、刷题训练、模拟考试、数据复盘、论坛互动与平台运营在一个系统中完整闭环。', 'size': 17, 'color': 'DDE7FF', 'space_after': 8},
                {'text': '谢谢各位老师指导', 'size': 14, 'color': 'B8C7DC'},
            ],
        )
        add_panel(slide, Inches(0.72), Inches(4.18), Inches(5.6), Inches(1.72), fill='18263F', line='31415F')
        add_rich_lines(
            slide,
            Inches(1.02), Inches(4.52), Inches(4.9), Inches(1.1),
            [
                {'text': '正式访问地址', 'size': 12.5, 'color': 'AFC2DD'},
                {'text': 'https://saksk.top', 'size': 20, 'color': 'FFFFFF', 'bold': True, 'font': EN_FONT_FAMILY},
                {'text': '现场网络受限时可切换本地 Docker 演示环境。', 'size': 11.8, 'color': 'B8C7DC'},
            ],
        )
        add_browser_frame(slide, Inches(7.15), Inches(1.08), Inches(5.56), Inches(3.14), SCREENSHOT_DIR / '10-刷题AI解析.png', '刷题 + AI解析', mode='contain', temp_dir=temp_dir)
        add_browser_frame(slide, Inches(7.52), Inches(4.36), Inches(3.42), Inches(2.12), SCREENSHOT_DIR / '15-后台论坛管理.png', '后台论坛管理', mode='contain', temp_dir=temp_dir)
        add_phone_frame(slide, Inches(11.04), Inches(4.02), Inches(1.78), Inches(2.46), SCREENSHOT_DIR / '09-首页-移动端.png', crop=(0, 0, 1, 0.23), temp_dir=temp_dir)
        add_footer(slide, 15, dark=True)

        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        prs.save(OUTPUT_PPTX)

    return OUTPUT_PPTX


if __name__ == '__main__':
    path = build_presentation()
    print(path)
