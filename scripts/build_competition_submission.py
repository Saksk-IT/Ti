from __future__ import annotations

import os
import subprocess
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt
import imageio.v2 as imageio
import numpy as np
from PIL import Image as PILImage
from PIL import ImageDraw, ImageFont
from pypdf import PdfReader
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, StyleSheet1, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from pptx import Presentation
from pptx.util import Inches


PROJECT_NAME = 'Sak-AI答题助手'
SUBMISSION_ID = '2025012345'
PROJECT_FOLDER_NAME = f'{SUBMISSION_ID}-参赛总文件夹'
PROJECT_NUMBER = SUBMISSION_ID
PUBLIC_URL = 'https://saksk.top'
LOCAL_URL = 'http://127.0.0.1:8000'
AUTHOR_1 = '王为硕'
AUTHOR_2 = '队员B'
TEACHER_1 = '指导教师（待填写）'
VERSION = 'V1.0'
DATE_TEXT = '2026年4月6日'

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_ROOT = ROOT / 'output' / 'doc' / PROJECT_FOLDER_NAME
DIR_01 = OUTPUT_ROOT / f'{SUBMISSION_ID}-01作品与答辩材料'
DIR_02 = OUTPUT_ROOT / f'{SUBMISSION_ID}-02素材与源码'
DIR_03 = OUTPUT_ROOT / f'{SUBMISSION_ID}-03设计与开发文档'
DIR_04 = OUTPUT_ROOT / f'{SUBMISSION_ID}-04作品演示视频'
DIR_IMG = DIR_02 / '插图'
SHOT_DIR = DIR_02 / '网站截图'
TMP_DIR = ROOT / 'tmp' / 'docs'
SCREENSHOT_SOURCE_DIR = ROOT / '.submission_assets' / 'screenshots'

SUMMARY_TEMPLATE = Path('/Users/saksk/Downloads/01 软件应用与开发作品提交要求/01-2 作品信息摘要模板（2025版）V2.docx')
SUMMARY_DOCX = DIR_03 / f'中国大学生计算机设计大赛作品信息概要表-{PROJECT_NAME}.docx'
SUMMARY_PDF = DIR_03 / f'中国大学生计算机设计大赛作品信息概要表-{PROJECT_NAME}.pdf'
DESIGN_DOCX = DIR_03 / f'软件应用与开发类作品设计和开发文档-{PROJECT_NAME}.docx'
DESIGN_PDF = DIR_03 / f'软件应用与开发类作品设计和开发文档-{PROJECT_NAME}.pdf'
RUNTIME_DOCX = DIR_01 / f'{PROJECT_NAME}-运行网址与答辩说明.docx'
RUNTIME_PDF = DIR_01 / f'{PROJECT_NAME}-运行网址与答辩说明.pdf'
PPT_DOCX = DIR_01 / f'{PROJECT_NAME}-答辩PPT提纲.docx'
PPT_PDF = DIR_01 / f'{PROJECT_NAME}-答辩PPT提纲.pdf'
PPTX_FILE = DIR_01 / f'{PROJECT_NAME}-答辩演示PPT.pptx'
PPTX_PREVIEW_PDF = DIR_01 / f'{PROJECT_NAME}-答辩演示PPT.pdf'
SPEECH_DOCX = DIR_01 / f'{PROJECT_NAME}-答辩讲稿要点.docx'
SPEECH_PDF = DIR_01 / f'{PROJECT_NAME}-答辩讲稿要点.pdf'
VIDEO_DOCX = DIR_04 / f'{PROJECT_NAME}-演示视频脚本.docx'
VIDEO_PDF = DIR_04 / f'{PROJECT_NAME}-演示视频脚本.pdf'
DEMO_VIDEO_MP4 = DIR_04 / f'{PROJECT_NAME}-作品运行演示视频.mp4'
SOURCE_ZIP = DIR_02 / f'{PROJECT_NAME}-源码包.zip'
WORK_URL_TXT = DIR_01 / f'{PROJECT_NAME}-访问地址.txt'
MATERIALS_TXT = DIR_02 / f'{PROJECT_NAME}-代表性素材说明.txt'
VIDEO_README = DIR_04 / 'README.txt'
ARCH_IMG = DIR_IMG / '系统架构图.png'
FLOW_IMG = DIR_IMG / '核心流程图.png'
PACKAGE_ZIP = ROOT / 'output' / 'doc' / f'{PROJECT_FOLDER_NAME}.zip'
TMP_PPT_SLIDES = TMP_DIR / 'ppt_slides'

FONT_CANDIDATES = [
    Path('/System/Library/Fonts/STHeiti Medium.ttc'),
    Path('/System/Library/Fonts/Hiragino Sans GB.ttc'),
    Path('/Library/Fonts/Arial Unicode.ttf'),
]

SCREENSHOTS = [
    ('01-首页.png', '/hub', '登录后的首页汇总题库入口、继续学习、签到与学习统计，是用户进入系统后的主控面板。'),
    ('02-题库广场.png', '/public/banks', '题库广场按系统题库与用户公开题库统一呈现，支持浏览、筛选与加入题库。'),
    ('03-题库名片详情.png', '/public/banks/card/user/46', '题库名片详情页展示题库简介、参与人数、加入方式与继续练习入口。'),
    ('04-我的题库.png', '/user/banks', '我的题库聚合用户创建与加入的题库资源，方便继续练习、整理和管理。'),
    ('05-数据中心.png', '/data', '数据中心通过趋势图、热力图和能力画像等可视化组件反馈学习状态。'),
    ('06-考试选择页.png', '/exams/select', '考试中心支持按公共题库或个人题库选择考试来源，进入模拟考试流程。'),
    ('07-论坛首页.png', '/forum', '论坛模块支持按版块浏览帖子、查看热帖并发起交流，强化学习社区互动。'),
    ('08-后台仪表盘.png', '/admin/dashboard', '后台仪表盘用于查看题目、科目、用户等系统级统计信息，支撑平台运营。'),
    ('09-首页-移动端.png', '/hub（移动端）', '移动端适配截图展示 Web 页面在窄屏设备下的响应式布局与触控友好交互。'),
]

FILE_ENTRIES = [
    (f'{SUBMISSION_ID}-01作品与答辩材料/运行网址与答辩说明', '运行网址、亮点总结、答辩顺序与截图说明，含本地部署地址与正式演示地址。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-01作品与答辩材料/答辩演示相关文档', '包含作品信息概要表、设计和开发文档、讲稿与视频脚本等正式说明材料。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-02素材与源码/网站截图与插图', '真实网站截图、系统架构图和核心流程图等创作素材。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-02素材与源码/源码包.zip', '提交用源码压缩包，包含后端、Web、小程序、Docker与文档等代表性源码。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-03设计与开发文档/答辩演示PPT', '实际答辩演示PPT，含首页、架构、核心功能、测试与总结等正式答辩页。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-03设计与开发文档/设计和开发文档', '作品信息概要表与软件应用开发文档等正式提交文档。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-04作品演示视频/作品运行演示视频.mp4', '基于真实网站截图制作的正式演示视频，覆盖首页、题库、考试、论坛与后台等环节。', '已上传到网盘', '自制'),
    (f'{SUBMISSION_ID}-04作品演示视频/视频脚本与说明', '补充视频脚本与目录说明，便于后续重录或替换。', '已上传到网盘', '自制'),
]


@dataclass
class SlideSection:
    title: str
    bullets: List[str]


INTRO_TEXT = 'Sak-AI答题助手是面向备考场景的Web题库平台，集题库广场、个人题库、刷题、考试、数据中心、论坛与后台管理于一体，并配套微信小程序延展移动学习场景。'
INNOVATION_TEXT = '作品采用Flask单仓全栈架构，实现Web与小程序共享语义；支持Docker一键部署、学习数据可视化与AI题目解析，兼顾工程化、扩展性和教学应用价值。'
SPECIAL_NOTE = (
    '1. 本作品不含涉及疆域展示的地图内容。\n'
    '2. 作品基于团队自研题库平台持续完善，本次参赛重点整理并强化 Web 主体演示链路，以及题库广场、个人题库、数据中心、论坛与后台等核心模块的闭环展示。\n'
    '3. 作品功能中集成 AI 题目解析能力，采用 DashScope 兼容接口生成题目解析文本；AI 仅作为局部辅助能力，不替代题库管理、练习、考试和数据统计等核心业务逻辑。'
)

COMPETITOR_ROWS = [
    ['对比维度', '常见刷题平台', f'{PROJECT_NAME}'],
    ['部署方式', '多为独立前端 + API 服务', '单仓全栈，Docker 统一部署'],
    ['题库结构', '偏公共题库或单一业务线', '公共题库与个人题库双线并存'],
    ['学习闭环', '练习与复盘分散', '练习、错题、收藏、考试、数据中心一体化'],
    ['运营能力', '社区与后台常需外挂系统', '内建论坛与后台管理模块'],
    ['移动场景', '通常需独立 App 或 H5 适配', '微信小程序复用同一后端语义'],
]

DATABASE_ROWS = [
    ['实体/对象', '作用说明'],
    ['用户（User）', '保存账户、角色、资料与登录态信息，支撑 Session/JWT 双认证。'],
    ['题库（公共/个人）', '区分系统题库、用户公开题库与个人题库，承载不同使用场景。'],
    ['题目（Question）', '保存题型、题干、选项、答案与解析等核心数据。'],
    ['学习记录', '记录答题结果、用户答案、做题时间、正确率与复盘数据。'],
    ['考试记录', '记录模拟考试的创建、提交、得分与结果统计。'],
    ['论坛内容', '记录帖子、评论、点赞与互动数据，形成学习社区。'],
]

TEST_ROWS = [
    ['测试项', '验证方法', '结果'],
    ['容器服务运行', '执行 docker compose --env-file .env -f compose.dev.yml ps', 'web/postgres/redis/worker 运行正常'],
    ['基础接口', f'访问 {LOCAL_URL}/api/ping', '返回 status=success'],
    ['深度健康检查', f'访问 {LOCAL_URL}/api/ping?deep=1', '返回 db=true、redis=true'],
    ['首页与导航', '浏览 /hub 并检查题库、考试、论坛入口', '页面可访问，导航正确'],
    ['题库广场', '浏览 /public/banks 与题库详情页', '列表、详情与题库信息展示正常'],
    ['个人题库', '浏览 /user/banks', '个人题库聚合展示正常'],
    ['数据中心', '浏览 /data', '图表容器与统计指标正常加载'],
    ['考试中心', '浏览 /exams/select', '题库选择与考试入口正常'],
    ['论坛与后台', '浏览 /forum、/admin/dashboard', '社区与后台统计页面可访问'],
]

PPT_SECTIONS = [
    SlideSection('封面', [f'作品名称：{PROJECT_NAME}', '赛道：软件应用与开发 - Web 应用与开发', f'正式访问地址：{PUBLIC_URL}']),
    SlideSection('问题背景', ['备考场景中题库分散、复盘弱、考试训练不足。', '教学与运营侧还需要统一管理与数据反馈能力。']),
    SlideSection('作品定位', ['以 Web 应用为主体，提供完整学习闭环。', '微信小程序作为移动延展端，共享同一后端与数据语义。']),
    SlideSection('整体架构', ['Flask 单仓全栈架构。', 'PostgreSQL + Redis + RQ + Docker 组成运行基础设施。']),
    SlideSection('核心功能 1', ['题库广场与题库名片详情。', '个人题库与分享加入。', '刷题、错题、收藏与用户进度。']),
    SlideSection('核心功能 2', ['模拟考试流程。', '数据中心可视化分析。', '论坛互动与后台管理。']),
    SlideSection('关键创新', ['双端共享语义，避免 Web 与小程序数据割裂。', 'Docker 统一部署，便于演示和复现。', 'AI 题目解析增强学习体验。']),
    SlideSection('真实运行展示', ['展示首页、题库广场、数据中心、考试、论坛、后台等真实页面截图。', '说明正式访问地址与本地演示地址。']),
    SlideSection('测试与验证', ['健康检查接口与容器服务验证通过。', '关键页面均已实际访问并截图留档。']),
    SlideSection('总结与展望', ['作品具备较完整的工程化形态与应用价值。', '后续可继续加强推荐、协作与数据分析能力。']),
]

SPEECH_SECTIONS = [
    SlideSection('第1页 封面', ['各位老师好，我们的作品是 Sak-AI答题助手，参赛类别为软件应用与开发中的 Web 应用与开发。', '本作品正式访问地址为 https://saksk.top，本地演示地址为 http://127.0.0.1:8000。']),
    SlideSection('第2页 问题背景', ['我们关注的是大学生备考与题库运营场景。', '传统工具常见问题包括题库割裂、错题复盘不足、学习数据不直观，以及教学运营端缺乏统一管理。']),
    SlideSection('第3页 作品定位', ['本作品以 Web 为主体，承载核心业务流程。', '同时配套微信小程序作为移动学习入口，但不另起一套数据语义。']),
    SlideSection('第4页 整体架构', ['系统采用 Flask 单仓全栈架构。', '后端连接 PostgreSQL 和 Redis，使用 RQ 处理异步任务，并通过 Docker 统一管理开发运行环境。']),
    SlideSection('第5页 核心功能一', ['在题库层面，我们提供题库广场、题库名片详情和个人题库。', '用户既可以浏览公共题库，也可以维护自己的私有与共享题库。']),
    SlideSection('第6页 核心功能二', ['在学习层面，系统支持刷题、错题、收藏、学习进度与模拟考试。', '在反馈层面，数据中心可以可视化呈现用户的正确率、覆盖率与学习节奏。']),
    SlideSection('第7页 核心功能三', ['系统还提供论坛模块支持交流互动。', '后台管理可以查看系统级统计，支持平台运营与维护。']),
    SlideSection('第8页 创新点', ['第一，Web 与小程序共享同一后端语义。', '第二，Docker 化部署提升了复现效率。', '第三，AI 题目解析增强了学习辅助能力。']),
    SlideSection('第9页 真实运行与测试', ['我们已经对首页、题库广场、个人题库、数据中心、考试、论坛和后台进行真实访问截图。', '同时通过健康检查接口与容器状态验证系统运行正常。']),
    SlideSection('第10页 总结', ['Sak-AI答题助手不仅是一个展示型网站，而是一个围绕备考场景构建的完整业务系统。', '谢谢各位老师，请批评指正。']),
]

VIDEO_SECTIONS = [
    SlideSection('片头 0:00-0:20', ['展示作品名称、赛道与正式访问地址 https://saksk.top。', '口播简要说明作品定位：面向备考学习与题库运营的 Web 应用。']),
    SlideSection('首页 0:20-1:00', ['打开 /hub，展示首页导航、学习入口、签到与学习概览。', '强调首页是用户进入系统后的统一工作台。']),
    SlideSection('题库广场 1:00-2:00', ['打开 /public/banks，展示系统题库与用户公开题库。', '说明题库广场支持浏览、筛选、加入与进入题库详情。']),
    SlideSection('题库名片详情 2:00-2:40', ['打开 /public/banks/card/user/46。', '说明题库简介、参与人数、加入方式与继续练习入口。']),
    SlideSection('我的题库 2:40-3:20', ['打开 /user/banks。', '说明“我创建的题库”“公开加入”“分享加入”的分类管理。']),
    SlideSection('数据中心 3:20-4:20', ['打开 /data。', '展示趋势图、热力图、能力画像与复盘建议。']),
    SlideSection('考试中心 4:20-5:00', ['打开 /exams/select。', '说明公共题库/个人题库考试入口与模拟考试流程。']),
    SlideSection('论坛首页 5:00-5:40', ['打开 /forum。', '展示版块、热帖与发帖入口，说明社区交流功能。']),
    SlideSection('后台仪表盘 5:40-6:20', ['打开 /admin/dashboard。', '展示题目数、科目数、用户数等后台统计能力。']),
    SlideSection('结尾 6:20-6:40', ['回到作品总结页，概括“题库-练习-考试-数据-社区-后台”的闭环。', '提示正式提交时将此脚本录制为不超过10分钟的 mp4 视频。']),
]

PPT_SLIDES = [
    {
        'title': PROJECT_NAME,
        'subtitle': '软件应用与开发 - Web 应用与开发\n竞赛答辩演示',
        'bullets': [
            f'正式访问地址：{PUBLIC_URL}',
            f'本地演示地址：{LOCAL_URL}',
            '单仓全栈题库平台，覆盖题库、练习、考试、数据中心、论坛与后台。',
        ],
        'image': SHOT_DIR / '01-首页.png',
    },
    {
        'title': '问题背景与作品定位',
        'subtitle': '围绕备考学习与题库运营场景构建完整 Web 应用',
        'bullets': [
            '备考工具常存在题库分散、错题复盘弱、考试训练不足的问题。',
            '本作品以 Web 为主体，小程序为移动延展端，共享同一套后端与数据语义。',
            '目标是把题库管理、练习、考试、数据反馈与社区互动打通成闭环。',
        ],
        'image': SHOT_DIR / '02-题库广场.png',
    },
    {
        'title': '系统架构',
        'subtitle': 'Flask + PostgreSQL + Redis + RQ + Docker',
        'bullets': [
            'Flask 应用工厂 + Blueprint 模块化组织。',
            'PostgreSQL 存储题库、考试与用户数据；Redis 支持缓存、限流与任务队列。',
            'Docker 开发模式便于快速部署、演示与复现。',
        ],
        'image': ARCH_IMG,
    },
    {
        'title': '题库体系',
        'subtitle': '公共题库 + 个人题库双线并行',
        'bullets': [
            '题库广场集中展示系统题库和用户公开题库。',
            '题库名片页支持查看简介、参与人数、加入方式与继续练习入口。',
            '个人题库支持用户整理自己的长期学习资源。',
        ],
        'image': SHOT_DIR / '03-题库名片详情.png',
    },
    {
        'title': '学习闭环',
        'subtitle': '题库浏览 → 练习/考试 → 数据复盘',
        'bullets': [
            '用户可从首页进入题库广场，也可从“我的题库”继续上次学习。',
            '系统支持错题、收藏、进度同步与考试训练。',
            '数据中心将练习结果转化为可视化反馈。',
        ],
        'image': FLOW_IMG,
    },
    {
        'title': '数据中心',
        'subtitle': '用图表驱动复盘与决策',
        'bullets': [
            '提供趋势图、热力图、能力画像、错题/收藏新增等视图。',
            '帮助用户从“做了多少题”升级到“知道自己哪里薄弱”。',
            '适合课程练习、阶段复盘与长期备考监控。',
        ],
        'image': SHOT_DIR / '05-数据中心.png',
    },
    {
        'title': '考试与社区',
        'subtitle': '训练结果与交流协同并存',
        'bullets': [
            '考试中心支持按公共题库或个人题库发起模拟考试。',
            '论坛模块支持版块浏览、热帖查看与学习交流。',
            '让训练、反馈和经验分享形成连续体验。',
        ],
        'image': SHOT_DIR / '07-论坛首页.png',
    },
    {
        'title': '后台与移动适配',
        'subtitle': '兼顾运营管理与多端体验',
        'bullets': [
            '后台仪表盘可查看题目数、科目数、用户数等系统统计信息。',
            'Web 页面支持移动端响应式布局，便于在窄屏设备上访问。',
            '小程序延展移动学习场景，但不复制一套独立数据逻辑。',
        ],
        'image': SHOT_DIR / '08-后台仪表盘.png',
    },
    {
        'title': '测试与验证',
        'subtitle': '真实运行、真实截图、真实接口检查',
        'bullets': [
            'docker compose 服务状态正常：web/postgres/redis/worker 全部运行。',
            '/api/ping 与 /api/ping?deep=1 已验证 success、db=true、redis=true。',
            '首页、题库广场、我的题库、数据中心、考试、论坛、后台均已实际访问截图。',
        ],
        'image': SHOT_DIR / '09-首页-移动端.png',
    },
    {
        'title': '总结与展望',
        'subtitle': '不仅是展示站，更是完整业务系统',
        'bullets': [
            '作品已形成题库-练习-考试-数据-社区-后台的一体化 Web 形态。',
            '工程化、可扩展、可部署、可演示，是本作品的核心优势。',
            '后续可继续增强个性化推荐、协作维护与更精细的数据分析能力。',
        ],
        'image': SHOT_DIR / '01-首页.png',
    },
]


def ensure_dirs() -> None:
    if OUTPUT_ROOT.exists():
        shutil.rmtree(OUTPUT_ROOT)
    if TMP_DIR.exists():
        shutil.rmtree(TMP_DIR)
    for path in [OUTPUT_ROOT, DIR_01, DIR_02, DIR_03, DIR_04, DIR_IMG, SHOT_DIR, TMP_DIR, TMP_PPT_SLIDES]:
        path.mkdir(parents=True, exist_ok=True)


def copy_submission_assets() -> None:
    missing: list[str] = []
    for filename, _, _ in SCREENSHOTS:
        src = SCREENSHOT_SOURCE_DIR / filename
        dst = SHOT_DIR / filename
        if not src.exists():
            missing.append(str(src))
            continue
        shutil.copy2(src, dst)
    if missing:
        raise FileNotFoundError('缺少截图源文件：\n' + '\n'.join(missing))


def set_cn_font(run, font_name: str = '微软雅黑', size: int | None = None, bold: bool | None = None) -> None:
    run.font.name = font_name
    r = run._element
    r.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def configure_doc(document: Document) -> None:
    sec = document.sections[0]
    sec.top_margin = Cm(2.2)
    sec.bottom_margin = Cm(2.0)
    sec.left_margin = Cm(2.4)
    sec.right_margin = Cm(2.2)
    normal = document.styles['Normal']
    normal.font.name = '微软雅黑'
    normal._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    normal.font.size = Pt(11)
    for style_name in ['Title', 'Heading 1', 'Heading 2', 'Heading 3']:
        if style_name in document.styles:
            style = document.styles[style_name]
            style.font.name = '微软雅黑'
            style._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')


def add_paragraph(document: Document, text: str, *, style: str | None = None, align=None, bold: bool = False, size: int = 11) -> None:
    p = document.add_paragraph(style=style)
    if align is not None:
        p.alignment = align
    for index, line in enumerate(text.split('\n')):
        if index:
            p.add_run('\n')
        run = p.add_run(line)
        set_cn_font(run, size=size, bold=bold)


def add_bullets(document: Document, items: Iterable[str], level: int = 0) -> None:
    for item in items:
        p = document.add_paragraph(style='List Bullet')
        if level:
            p.paragraph_format.left_indent = Cm(0.74 * level)
        run = p.add_run(item)
        set_cn_font(run, size=11)


def add_numbered(document: Document, items: Iterable[str]) -> None:
    for item in items:
        p = document.add_paragraph(style='List Number')
        run = p.add_run(item)
        set_cn_font(run, size=11)


def set_cell_text(cell, text: str, font_size: int = 10, bold: bool = False) -> None:
    cell.text = ''
    lines = text.split('\n') if text else ['']
    for i, line in enumerate(lines):
        p = cell.paragraphs[0] if i == 0 else cell.add_paragraph()
        p.paragraph_format.space_after = Pt(0)
        run = p.add_run(line)
        set_cn_font(run, size=font_size, bold=bold)


def set_paragraph_text_preserve(paragraph, text: str, font_name: str = '仿宋') -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ''
        paragraph.runs[0].font.name = font_name
        paragraph.runs[0]._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    else:
        run = paragraph.add_run(text)
        run.font.name = font_name
        run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)


def set_cell_paragraph(cell, paragraph_index: int, text: str, font_name: str = '仿宋') -> None:
    while len(cell.paragraphs) <= paragraph_index:
        cell.add_paragraph()
    set_paragraph_text_preserve(cell.paragraphs[paragraph_index], text, font_name=font_name)


def export_docx_to_pdf_with_pages(docx_path: Path, pdf_path: Path) -> None:
    script_path = TMP_DIR / 'export_pages_pdf.applescript'
    script_path.write_text(
        'on run argv\n'
        '  set inputPath to POSIX file (item 1 of argv)\n'
        '  set outputPath to POSIX file (item 2 of argv)\n'
        '  tell application "Pages"\n'
        '    activate\n'
        '    set theDoc to open inputPath\n'
        '    export theDoc to outputPath as PDF\n'
        '    close theDoc saving no\n'
        '    quit\n'
        '  end tell\n'
        'end run\n',
        encoding='utf-8',
    )
    subprocess.run(
        ['osascript', str(script_path), str(docx_path), str(pdf_path)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def add_table(document: Document, rows: List[List[str]], widths_cm: List[float] | None = None, font_size: int = 10) -> None:
    table = document.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = 'Table Grid'
    table.autofit = False
    if widths_cm:
        for row in table.rows:
            for cell, width in zip(row.cells, widths_cm):
                cell.width = Cm(width)
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            set_cell_text(table.cell(r_idx, c_idx), value, font_size=font_size, bold=(r_idx == 0))


def insert_image_after_paragraph(paragraph, image_path: Path, width_cm: float = 15.8):
    new_p = OxmlElement('w:p')
    paragraph._p.addnext(new_p)
    new_para = paragraph._parent.add_paragraph()
    new_para._p = new_p
    run = new_para.add_run()
    run.add_picture(str(image_path), width=Cm(width_cm))
    new_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    return new_para


def get_font_path() -> str:
    for path in FONT_CANDIDATES:
        if path.exists():
            return str(path)
    raise FileNotFoundError('未找到可用于绘图的中文字体文件。')


def create_diagram(path: Path, title: str, boxes: List[tuple[int, int, int, int, str]], arrows: List[tuple[tuple[int, int], tuple[int, int]]], footer: str) -> None:
    font_path = get_font_path()
    img = PILImage.new('RGB', (1500, 900), 'white')
    draw = ImageDraw.Draw(img)
    title_font = ImageFont.truetype(font_path, 40)
    box_font = ImageFont.truetype(font_path, 26)
    footer_font = ImageFont.truetype(font_path, 22)
    draw.text((750, 48), title, anchor='mm', fill='#0F172A', font=title_font)
    for x1, y1, x2, y2, text in boxes:
        draw.rounded_rectangle((x1, y1, x2, y2), radius=24, fill='#F8FAFC', outline='#2563EB', width=4)
        draw.multiline_text(((x1 + x2) / 2, (y1 + y2) / 2), text, anchor='mm', fill='#0F172A', font=box_font, align='center', spacing=8)
    for start, end in arrows:
        draw.line([start, end], fill='#475569', width=6)
        ex, ey = end
        sx, sy = start
        if abs(ex - sx) >= abs(ey - sy):
            direction = 1 if ex >= sx else -1
            draw.polygon([(ex, ey), (ex - 20 * direction, ey - 12), (ex - 20 * direction, ey + 12)], fill='#475569')
        else:
            direction = 1 if ey >= sy else -1
            draw.polygon([(ex, ey), (ex - 12, ey - 20 * direction), (ex + 12, ey - 20 * direction)], fill='#475569')
    draw.text((750, 850), footer, anchor='mm', fill='#475569', font=footer_font)
    img.save(path)


def build_diagrams() -> None:
    create_diagram(
        ARCH_IMG,
        'Sak-AI答题助手系统架构图',
        [
            (100, 200, 420, 330, 'Web 浏览器'),
            (100, 420, 420, 550, '微信小程序'),
            (560, 310, 980, 470, 'Flask Web / API\n应用工厂 + Blueprint'),
            (1100, 160, 1400, 280, 'PostgreSQL 16\n题库/用户/考试数据'),
            (1100, 360, 1400, 480, 'Redis 7\n缓存/限流/RQ'),
            (1100, 560, 1400, 680, 'RQ Worker\n异步任务'),
            (560, 590, 980, 720, 'Docker 开发环境\ncompose.dev.yml'),
        ],
        [
            ((420, 265), (560, 365)),
            ((420, 485), (560, 415)),
            ((980, 360), (1100, 220)),
            ((980, 390), (1100, 420)),
            ((980, 665), (1100, 620)),
            ((1250, 560), (1250, 480)),
        ],
        'Web 为主应用端，小程序复用同一后端与数据语义。',
    )
    create_diagram(
        FLOW_IMG,
        '核心用户流程图',
        [
            (80, 300, 300, 420, '首页 /hub'),
            (360, 300, 620, 420, '题库广场 /public/banks'),
            (680, 300, 960, 420, '题库详情 / 我的题库'),
            (1020, 180, 1360, 300, '刷题练习 / 题目详情'),
            (1020, 420, 1360, 540, '模拟考试 /exams/select'),
            (1020, 660, 1360, 780, '数据复盘 /data'),
        ],
        [
            ((300, 360), (360, 360)),
            ((620, 360), (680, 360)),
            ((960, 330), (1020, 240)),
            ((960, 390), (1020, 480)),
            ((1190, 540), (1190, 660)),
        ],
        '从题库浏览到练习、考试与复盘，形成完整学习闭环。',
    )


def wrap_text_for_pil(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = list(text)
    lines: list[str] = []
    current = ''
    for ch in words:
        candidate = current + ch
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = ch
    if current:
        lines.append(current)
    return lines


def fit_image(image: PILImage.Image, max_width: int, max_height: int) -> PILImage.Image:
    ratio = min(max_width / image.width, max_height / image.height)
    new_size = (max(1, int(image.width * ratio)), max(1, int(image.height * ratio)))
    return image.resize(new_size, PILImage.Resampling.LANCZOS)


def render_ppt_slide_image(slide_data: dict, index: int, total: int, output_path: Path) -> None:
    width, height = 1920, 1080
    img = PILImage.new('RGB', (width, height), '#F8FAFC')
    draw = ImageDraw.Draw(img)
    font_path = get_font_path()
    title_font = ImageFont.truetype(font_path, 52)
    subtitle_font = ImageFont.truetype(font_path, 28)
    bullet_font = ImageFont.truetype(font_path, 30)
    small_font = ImageFont.truetype(font_path, 22)

    draw.rectangle((0, 0, width, 112), fill='#0F172A')
    draw.text((84, 58), slide_data['title'], fill='white', font=title_font, anchor='lm')
    draw.text((1770, 58), f'{index:02d}/{total:02d}', fill='#CBD5E1', font=small_font, anchor='rm')

    content_x = 84
    content_y = 156
    draw.text((content_x, content_y), slide_data.get('subtitle', ''), fill='#334155', font=subtitle_font)

    bullet_y = content_y + 82
    bullet_area_width = 700
    for bullet in slide_data.get('bullets', []):
        lines = wrap_text_for_pil(draw, bullet, bullet_font, bullet_area_width)
        draw.ellipse((content_x, bullet_y + 8, content_x + 14, bullet_y + 22), fill='#2563EB')
        text_x = content_x + 30
        for line in lines:
            draw.text((text_x, bullet_y), line, fill='#0F172A', font=bullet_font)
            bullet_y += 40
        bullet_y += 14

    image_box = (860, 170, 1820, 860)
    draw.rounded_rectangle(image_box, radius=28, fill='white', outline='#CBD5E1', width=3)
    image_path = Path(slide_data['image'])
    shot = PILImage.open(image_path).convert('RGB')
    shot_fit = fit_image(shot, image_box[2] - image_box[0] - 28, image_box[3] - image_box[1] - 28)
    paste_x = image_box[0] + ((image_box[2] - image_box[0] - shot_fit.width) // 2)
    paste_y = image_box[1] + ((image_box[3] - image_box[1] - shot_fit.height) // 2)
    img.paste(shot_fit, (paste_x, paste_y))

    caption_box = (80, 900, 1840, 1020)
    draw.rounded_rectangle(caption_box, radius=24, fill='#E2E8F0')
    caption = slide_data.get('caption') or ''
    draw.text((110, 948), caption, fill='#334155', font=small_font)
    draw.text((1810, 948), PROJECT_NAME, fill='#475569', font=small_font, anchor='rm')

    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(output_path)


def build_ppt_slide_images() -> list[Path]:
    for old_png in TMP_PPT_SLIDES.glob('slide-*.png'):
        old_png.unlink()
    slides: list[Path] = []
    total = len(PPT_SLIDES)
    for idx, slide in enumerate(PPT_SLIDES, start=1):
        out = TMP_PPT_SLIDES / f'slide-{idx:02d}.png'
        caption = slide.get('subtitle', '')
        slide = {**slide, 'caption': caption}
        render_ppt_slide_image(slide, idx, total, out)
        slides.append(out)
    return slides


def build_answer_ppt(slide_images: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for slide_img in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(slide_img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(PPTX_FILE))


def build_answer_ppt_pdf(slide_images: list[Path]) -> None:
    page_w = 13.333 * 72
    page_h = 7.5 * 72
    pdf = canvas.Canvas(str(PPTX_PREVIEW_PDF), pagesize=(page_w, page_h))
    for slide_img in slide_images:
        pdf.drawImage(str(slide_img), 0, 0, width=page_w, height=page_h, preserveAspectRatio=True, mask='auto')
        pdf.showPage()
    pdf.save()


def render_video_frame(title: str, route: str, caption: str, image_path: Path) -> PILImage.Image:
    width, height = 1280, 720
    canvas_img = PILImage.new('RGB', (width, height), '#0F172A')
    draw = ImageDraw.Draw(canvas_img)
    font_path = get_font_path()
    title_font = ImageFont.truetype(font_path, 36)
    text_font = ImageFont.truetype(font_path, 22)

    draw.rectangle((0, 0, width, 80), fill='#111827')
    draw.text((40, 40), title, fill='white', font=title_font, anchor='lm')
    draw.text((1240, 40), PUBLIC_URL, fill='#CBD5E1', font=text_font, anchor='rm')

    shot = PILImage.open(image_path).convert('RGB')
    shot_fit = fit_image(shot, 1180, 500)
    shot_x = (width - shot_fit.width) // 2
    shot_y = 110 + (500 - shot_fit.height) // 2
    canvas_img.paste(shot_fit, (shot_x, shot_y))

    draw.rounded_rectangle((28, 630, 1252, 700), radius=20, fill='#E2E8F0')
    draw.text((50, 654), f'页面路径：{route}', fill='#0F172A', font=text_font)
    draw.text((50, 682), caption, fill='#334155', font=text_font)
    return canvas_img


def render_title_frame(title: str, subtitle: str) -> PILImage.Image:
    width, height = 1280, 720
    img = PILImage.new('RGB', (width, height), '#0F172A')
    draw = ImageDraw.Draw(img)
    font_path = get_font_path()
    title_font = ImageFont.truetype(font_path, 58)
    sub_font = ImageFont.truetype(font_path, 28)
    draw.text((width / 2, 280), title, fill='white', font=title_font, anchor='mm')
    draw.text((width / 2, 360), subtitle, fill='#CBD5E1', font=sub_font, anchor='mm')
    draw.text((width / 2, 430), PUBLIC_URL, fill='#93C5FD', font=sub_font, anchor='mm')
    return img


def build_demo_video() -> None:
    fps = 24
    writer = imageio.get_writer(str(DEMO_VIDEO_MP4), fps=fps)
    intro = render_title_frame(PROJECT_NAME, '作品运行演示视频')
    outro = render_title_frame('演示结束', '题库-练习-考试-数据-社区-后台')

    for _ in range(fps * 2):
        writer.append_data(np.asarray(intro))

    for filename, route, caption in SCREENSHOTS:
        title = filename.replace('.png', '')
        frame = render_video_frame(title, route, caption, SHOT_DIR / filename)
        for _ in range(fps * 3):
            writer.append_data(np.asarray(frame))

    for _ in range(fps * 2):
        writer.append_data(np.asarray(outro))
    writer.close()


def build_package_zip() -> None:
    if PACKAGE_ZIP.exists():
        PACKAGE_ZIP.unlink()
    with zipfile.ZipFile(PACKAGE_ZIP, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=8) as zf:
        for file in OUTPUT_ROOT.rglob('*'):
            if file.is_dir() or file == PACKAGE_ZIP or file.name == '.DS_Store':
                continue
            arcname = Path(PROJECT_FOLDER_NAME) / file.relative_to(OUTPUT_ROOT)
            zf.write(file, arcname.as_posix())


def reportlab_styles() -> StyleSheet1:
    pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='CNTitle', parent=styles['Title'], fontName='STSong-Light', alignment=TA_CENTER, fontSize=22, leading=30, textColor=colors.HexColor('#0F172A')))
    styles.add(ParagraphStyle(name='CNSubTitle', parent=styles['Normal'], fontName='STSong-Light', alignment=TA_CENTER, fontSize=11, leading=16, textColor=colors.HexColor('#475569')))
    styles.add(ParagraphStyle(name='CNH1', parent=styles['Heading1'], fontName='STSong-Light', fontSize=16, leading=24, spaceBefore=10, spaceAfter=8, textColor=colors.HexColor('#0F172A')))
    styles.add(ParagraphStyle(name='CNH2', parent=styles['Heading2'], fontName='STSong-Light', fontSize=13, leading=20, spaceBefore=8, spaceAfter=6, textColor=colors.HexColor('#0F172A')))
    styles.add(ParagraphStyle(name='CNBody', parent=styles['BodyText'], fontName='STSong-Light', fontSize=10.5, leading=17, alignment=TA_JUSTIFY, textColor=colors.HexColor('#111827')))
    styles.add(ParagraphStyle(name='CNBodyLeft', parent=styles['BodyText'], fontName='STSong-Light', fontSize=10.5, leading=17, alignment=TA_LEFT, textColor=colors.HexColor('#111827')))
    styles.add(ParagraphStyle(name='CNCaption', parent=styles['BodyText'], fontName='STSong-Light', fontSize=9.5, leading=14, alignment=TA_CENTER, textColor=colors.HexColor('#475569')))
    return styles


def build_pdf(path: Path, title: str, subtitle_lines: List[str], body_builder) -> None:
    styles = reportlab_styles()
    doc = SimpleDocTemplate(str(path), pagesize=A4, leftMargin=2.1 * cm, rightMargin=2.1 * cm, topMargin=1.8 * cm, bottomMargin=1.8 * cm)
    story = [Paragraph(title, styles['CNTitle'])]
    for line in subtitle_lines:
        story.append(Paragraph(line, styles['CNSubTitle']))
    story.append(Spacer(1, 0.35 * cm))
    story.extend(body_builder(styles))
    doc.build(story)


def make_table_pdf(rows: List[List[str]], widths: List[float]) -> Table:
    table = Table(rows, colWidths=[w * cm for w in widths], repeatRows=1)
    table.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (-1, -1), 'STSong-Light'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('LEADING', (0, 0), (-1, -1), 12),
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#E8EEF9')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#0F172A')),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#94A3B8')),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
        ('LEFTPADDING', (0, 0), (-1, -1), 5),
        ('RIGHTPADDING', (0, 0), (-1, -1), 5),
        ('TOPPADDING', (0, 0), (-1, -1), 5),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
    ]))
    return table


def add_pdf_bullets(story: list, styles: StyleSheet1, items: Iterable[str]) -> None:
    for item in items:
        story.append(Paragraph(f'• {item}', styles['CNBodyLeft']))


def add_pdf_image(story: list, styles: StyleSheet1, image_path: Path, caption: str, max_width_cm: float = 15.5) -> None:
    pil = PILImage.open(image_path)
    width, height = pil.size
    ratio = height / width
    target_width = max_width_cm * cm
    target_height = target_width * ratio
    if target_height > 20 * cm:
        target_height = 20 * cm
        target_width = target_height / ratio
    story.append(Image(str(image_path), width=target_width, height=target_height))
    story.append(Spacer(1, 0.15 * cm))
    story.append(Paragraph(caption, styles['CNCaption']))
    story.append(Spacer(1, 0.25 * cm))


def build_summary_docx() -> None:
    doc = Document(str(SUMMARY_TEMPLATE))
    table = doc.tables[0]

    set_cell_paragraph(table.cell(0, 2), 0, '待填写')
    set_cell_paragraph(table.cell(0, 6), 0, PROJECT_NAME)
    set_cell_paragraph(table.cell(1, 2), 0, '软件应用与开发')
    set_cell_paragraph(table.cell(1, 9), 0, 'Web应用与开发')

    set_cell_paragraph(table.cell(2, 0), 1, 'Sak-AI答题助手是面向备考场景的Web题库平台，集题库广场、个人题库、刷题、考试、数据中心、论坛与后台管理于一体，')
    set_cell_paragraph(table.cell(2, 0), 2, '并配套微信小程序延展移动学习场景。')

    set_cell_paragraph(table.cell(3, 0), 1, '作品采用Flask单仓全栈架构，实现Web与小程序共享语义；支持Docker一键部署、学习数据可视化与')
    set_cell_paragraph(table.cell(3, 0), 2, 'AI题目解析，兼顾工程化、扩展性和教学应用价值。')

    set_cell_paragraph(table.cell(4, 0), 1, '1.本作品不含涉及疆域展示的地图内容。')
    set_cell_paragraph(table.cell(4, 0), 2, '2.作品基于团队自研平台持续完善，本次参赛重点整理并强化Web主体演示链路。')
    set_cell_paragraph(table.cell(4, 0), 3, '3.作品功能中集成AI题目解析能力，采用DashScope兼容接口生成解析文本。')

    set_cell_paragraph(table.cell(6, 2), 0, AUTHOR_1)
    set_cell_paragraph(table.cell(6, 5), 0, AUTHOR_2)
    set_cell_paragraph(table.cell(6, 8), 0, '')
    set_cell_paragraph(table.cell(6, 10), 0, '')
    set_cell_paragraph(table.cell(6, 12), 0, '')

    percentages = [('90%', '10%')] * 7
    for idx, (p1, p2) in enumerate(percentages, start=7):
        set_cell_paragraph(table.cell(idx, 2), 0, p1)
        set_cell_paragraph(table.cell(idx, 5), 0, p2)
        set_cell_paragraph(table.cell(idx, 8), 0, '')
        set_cell_paragraph(table.cell(idx, 10), 0, '')
        set_cell_paragraph(table.cell(idx, 12), 0, '')

    set_cell_paragraph(table.cell(15, 3), 0, '□作品创意 ■理论指导 ■技术方案 □实验场地 □硬件资源')
    set_cell_paragraph(table.cell(15, 3), 1, '□数据提供 □后勤支持 □宣讲通知 □组织协调 □经费支持')
    set_cell_paragraph(table.cell(15, 3), 2, f'■其他：姓名待补（{TEACHER_1}）')
    set_cell_paragraph(table.cell(16, 3), 0, '□Windows □Linux ■macOS □其他：')
    set_cell_paragraph(table.cell(17, 3), 0, '□Windows □Linux ■macOS □iOS □Android □其他：')
    set_cell_paragraph(table.cell(18, 3), 0, 'Python 3、Flask 3、SQLAlchemy 2、PostgreSQL 16、Redis 7、Docker、Jinja2、原生JS/CSS、微信开发者工具、Chrome DevTools、Git')
    set_cell_paragraph(table.cell(19, 3), 0, '1、Flask 官方文档')
    set_cell_paragraph(table.cell(19, 3), 1, '2、SQLAlchemy 官方文档')
    set_cell_paragraph(table.cell(19, 3), 2, '3、Docker 官方文档')
    set_cell_paragraph(table.cell(20, 3), 0, '■素材压缩包 ■报告文档 ■演示视频 ■PPT ■源代码 ■部署文件')
    set_cell_paragraph(table.cell(20, 3), 1, '□数据集 □模型 ■作品文件 □其他')

    for i, (name, desc, status, copyright_text) in enumerate(FILE_ENTRIES, start=23):
        set_cell_paragraph(table.cell(i, 1), 0, f'文件名：{name}')
        set_cell_paragraph(table.cell(i, 1), 1, f'描述：{desc}')
        set_cell_paragraph(table.cell(i, 7), 0, '■已上传到网盘')
        set_cell_paragraph(table.cell(i, 7), 1, '□未上传，下载地址：')
        set_cell_paragraph(table.cell(i, 11), 0, '■自制  □未知版权')
        set_cell_paragraph(table.cell(i, 11), 1, '□开源  □获得授权')

    set_cell_paragraph(table.cell(31, 0), 3, '全体参赛队员签名：（可附授权使用的电子签名图片）')
    doc.save(str(SUMMARY_DOCX))


def build_summary_pdf() -> None:
    export_docx_to_pdf_with_pages(SUMMARY_DOCX, SUMMARY_PDF)


def build_design_docx() -> None:
    doc = Document()
    configure_doc(doc)
    add_paragraph(doc, '中国大学生计算机设计大赛', style='Title', align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, size=18)
    add_paragraph(doc, '软件开发类作品设计和开发文档', align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, size=16)
    add_paragraph(doc, f'作品编号：{PROJECT_NUMBER}\n作品名称：{PROJECT_NAME}\n作　　者：{AUTHOR_1}、{AUTHOR_2}\n版本编号：{VERSION}\n填写日期：{DATE_TEXT}', align=WD_ALIGN_PARAGRAPH.CENTER, size=11)

    doc.add_heading('需求分析', level=1)
    add_paragraph(doc, 'Sak-AI答题助手面向大学生备考、课程练习与题库运营场景，解决题库分散、错题复盘弱、考试训练不足和学习数据不可视等问题。传统刷题工具往往偏重单点练习，公共题库与个人题库相互割裂，导致用户难以沉淀个人学习资源，也不利于教师或运营侧开展内容管理。本作品以 Web 应用为主体，通过统一的题库管理、练习、考试、数据分析、论坛互动和后台管理能力，为学习者和题库维护者提供完整闭环。')
    add_paragraph(doc, '目标用户主要包括三类：一是需要高频刷题和复盘的大学生与备考学习者；二是需要组织题库、审核内容和管理用户的教师或运营人员；三是希望随时随地访问同一套题库语义的移动端用户。对标常见在线刷题平台和通用题库系统，本作品强调工程化部署、共享数据语义以及学习闭环的一体化表达。')
    add_table(doc, COMPETITOR_ROWS, widths_cm=[3.0, 6.0, 7.0], font_size=10)

    doc.add_heading('概要设计', level=1)
    add_paragraph(doc, '系统整体采用 Flask 单仓全栈架构：Web 浏览器与微信小程序共同访问 Flask 提供的页面与 API 服务，后端连接 PostgreSQL 存储题库与学习数据，使用 Redis 承担缓存、限流与异步任务队列，RQ Worker 负责后台任务处理；本地开发通过 compose.dev.yml 启动完整环境。')
    add_paragraph(doc, '核心模块按照业务边界拆分为认证、主页面、题库练习、考试、用户、聊天、通知、弹窗、编程、个人题库、论坛与后台管理等 12 个 Blueprint 模块，兼顾功能内聚与后期扩展。')
    add_paragraph(doc, '模块划分如下：')
    add_bullets(doc, [
        '用户认证：登录、会话、JWT、角色权限。',
        '公共题库：题库广场、题库详情、题目浏览。',
        '个人题库：自建题库、公开分享、加入题库与管理。',
        '刷题与错题收藏：答题记录、错题、收藏、进度同步。',
        '模拟考试：考试创建、选择题库、提交与结果查看。',
        '数据中心：正确率、覆盖率、趋势图、热力图与能力画像。',
        '论坛社区：帖子、评论、版块与互动。',
        '后台管理：题目、用户、科目和系统统计管理。',
    ])
    p = doc.add_paragraph()
    r = p.add_run('图 1  系统整体架构图')
    set_cn_font(r, size=11, bold=True)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    insert_image_after_paragraph(p, ARCH_IMG, width_cm=15.8)

    doc.add_heading('详细设计', level=1)
    add_paragraph(doc, '用户在首页进入系统后，可以从题库广场浏览公共题库，也可以在“我的题库”中访问已创建或已加入的题库资源；随后进入练习与考试流程，在答题后由数据中心进行统计与复盘，形成连续的学习闭环。')
    p = doc.add_paragraph()
    r = p.add_run('图 2  核心用户流程图')
    set_cn_font(r, size=11, bold=True)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    insert_image_after_paragraph(p, FLOW_IMG, width_cm=15.8)
    add_paragraph(doc, '数据库设计以用户、题库、题目、学习记录、考试记录和论坛内容为核心实体，通过统一的数据语义支撑 Web 与小程序共用同一套后端逻辑。')
    add_table(doc, DATABASE_ROWS, widths_cm=[4.0, 12.0], font_size=10)
    add_paragraph(doc, '关键技术点包括：')
    add_numbered(doc, [
        '采用 Flask 应用工厂 + Blueprint 组织方式，便于按模块扩展。',
        '兼容 Session 与 JWT 双认证模式，分别服务 Web 与 API/小程序。',
        '接口返回逐步统一为 status/code/data/message 信封结构，降低前后端兼容成本。',
        '基于 Docker 的本地开发模式统一了 Web、Worker、PostgreSQL 与 Redis 运行环境。',
        '通过数据中心可视化组件将练习、考试、复盘转化为可感知的学习反馈。',
        'AI 题目解析能力通过兼容接口接入，用于增强题目讲解与学习辅助体验。',
    ])
    doc.add_heading('真实运行界面截图', level=2)
    for idx, (filename, route, caption) in enumerate(SCREENSHOTS, start=1):
        img_path = SHOT_DIR / filename
        add_paragraph(doc, f'图 {idx + 2}  {filename.replace(".png", "")}（路径：{route}）', align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, size=11)
        doc.add_picture(str(img_path), width=Cm(15.8))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
        add_paragraph(doc, caption, align=WD_ALIGN_PARAGRAPH.CENTER, size=10)

    doc.add_heading('测试报告', level=1)
    add_paragraph(doc, '作品运行环境采用 Docker 开发模式，服务由 compose.dev.yml 管理。验证时重点关注容器服务状态、健康检查接口与关键业务页面是否可访问。')
    add_table(doc, TEST_ROWS, widths_cm=[3.5, 6.5, 6.0], font_size=9)
    add_paragraph(doc, '验证截图均来自真实运行站点，正式访问地址为 https://saksk.top，本地演示地址为 http://127.0.0.1:8000。两者共用相同系统语义，其中正式地址用于公开访问，本地地址用于开发与现场演示说明。')

    doc.add_heading('安装及使用', level=1)
    add_numbered(doc, [
        '准备 .env 文件，并根据环境填写必要的配置项。',
        '在项目根目录执行：docker compose --env-file .env -f compose.dev.yml up。',
        '等待 web、postgres、redis、worker 容器启动完成。',
        f'浏览器访问 {LOCAL_URL} 进行本地演示。',
        f'正式展示可使用公网地址 {PUBLIC_URL}。',
        '根据页面导航依次演示首页、题库广场、题库详情、我的题库、数据中心、考试中心、论坛和后台。',
    ])

    doc.add_heading('项目总结', level=1)
    add_paragraph(doc, 'Sak-AI答题助手围绕“题库资源管理 + 学习闭环反馈 + 社区互动 + 运营支撑”构建完整 Web 应用，不仅实现了刷题功能，还将题库组织、考试训练、数据可视化、论坛交流和后台管理整合到同一系统中。工程上，作品采用单仓全栈架构，保证 Web 与小程序共享语义，降低维护成本；部署上，Docker 化运行方式提升了演示和复现效率；应用上，AI 题目解析能力让系统更适合面向教学与学习场景的持续扩展。后续可进一步增强推荐策略、协作编辑与个性化复盘能力。')

    doc.add_heading('参考文献', level=1)
    add_numbered(doc, [
        'Flask Official Documentation. https://flask.palletsprojects.com/',
        'SQLAlchemy Documentation. https://docs.sqlalchemy.org/',
        'Docker Docs. https://docs.docker.com/',
        '微信开放文档 - 小程序。https://developers.weixin.qq.com/miniprogram/dev/framework/',
    ])
    doc.save(str(DESIGN_DOCX))


def build_design_pdf() -> None:
    def body(styles: StyleSheet1):
        story = []
        story.append(Paragraph('作品编号：待填写<br/>作品名称：Sak-AI答题助手<br/>作者：王为硕、队员B（占位）<br/>版本编号：V1.0<br/>填写日期：2026年4月6日', styles['CNSubTitle']))
        story.append(Spacer(1, 0.25 * cm))

        story.append(Paragraph('1. 需求分析', styles['CNH1']))
        story.append(Paragraph('Sak-AI答题助手面向大学生备考、课程练习与题库运营场景，解决题库分散、错题复盘弱、考试训练不足和学习数据不可视等问题。传统刷题工具往往偏重单点练习，公共题库与个人题库相互割裂，导致用户难以沉淀个人学习资源，也不利于教师或运营侧开展内容管理。本作品以 Web 应用为主体，通过统一的题库管理、练习、考试、数据分析、论坛互动和后台管理能力，为学习者和题库维护者提供完整闭环。', styles['CNBody']))
        story.append(Paragraph('目标用户包括备考学习者、题库维护者和移动端用户。作品对标常见在线刷题平台和通用题库系统，但更强调工程化部署、共享数据语义以及学习闭环的一体化表达。', styles['CNBody']))
        story.append(make_table_pdf(COMPETITOR_ROWS, [3.2, 6.0, 6.0]))
        story.append(Spacer(1, 0.2 * cm))

        story.append(Paragraph('2. 概要设计', styles['CNH1']))
        story.append(Paragraph('系统整体采用 Flask 单仓全栈架构：Web 浏览器与微信小程序共同访问 Flask 提供的页面与 API 服务，后端连接 PostgreSQL 存储题库与学习数据，使用 Redis 承担缓存、限流与异步任务队列，RQ Worker 负责后台任务处理；本地开发通过 compose.dev.yml 启动完整环境。', styles['CNBody']))
        story.append(Paragraph('核心模块按照业务边界拆分为认证、主页面、题库练习、考试、用户、聊天、通知、弹窗、编程、个人题库、论坛与后台管理等 12 个 Blueprint 模块，兼顾功能内聚与后期扩展。', styles['CNBody']))
        add_pdf_bullets(story, styles, [
            '用户认证：登录、会话、JWT、角色权限。',
            '公共题库：题库广场、题库详情、题目浏览。',
            '个人题库：自建题库、公开分享、加入题库与管理。',
            '刷题与错题收藏：答题记录、错题、收藏、进度同步。',
            '模拟考试：考试创建、选择题库、提交与结果查看。',
            '数据中心：正确率、覆盖率、趋势图、热力图与能力画像。',
            '论坛社区：帖子、评论、版块与互动。',
            '后台管理：题目、用户、科目和系统统计管理。',
        ])
        add_pdf_image(story, styles, ARCH_IMG, '图 1  系统整体架构图')

        story.append(Paragraph('3. 详细设计', styles['CNH1']))
        story.append(Paragraph('用户在首页进入系统后，可以从题库广场浏览公共题库，也可以在“我的题库”中访问已创建或已加入的题库资源；随后进入练习与考试流程，在答题后由数据中心进行统计与复盘，形成连续的学习闭环。', styles['CNBody']))
        add_pdf_image(story, styles, FLOW_IMG, '图 2  核心用户流程图')
        story.append(Paragraph('数据库设计以用户、题库、题目、学习记录、考试记录和论坛内容为核心实体，通过统一的数据语义支撑 Web 与小程序共用同一套后端逻辑。', styles['CNBody']))
        story.append(make_table_pdf(DATABASE_ROWS, [4.0, 12.0]))
        story.append(Spacer(1, 0.2 * cm))
        add_pdf_bullets(story, styles, [
            '采用 Flask 应用工厂 + Blueprint 组织方式，便于按模块扩展。',
            '兼容 Session 与 JWT 双认证模式，分别服务 Web 与 API/小程序。',
            '接口返回逐步统一为 status/code/data/message 信封结构，降低前后端兼容成本。',
            '基于 Docker 的本地开发模式统一了 Web、Worker、PostgreSQL 与 Redis 运行环境。',
            '通过数据中心可视化组件将练习、考试、复盘转化为可感知的学习反馈。',
            'AI 题目解析能力通过兼容接口接入，用于增强题目讲解与学习辅助体验。',
        ])
        story.append(PageBreak())
        story.append(Paragraph('3.1 真实运行界面截图', styles['CNH2']))
        for idx, (filename, route, caption) in enumerate(SCREENSHOTS, start=1):
            add_pdf_image(story, styles, SHOT_DIR / filename, f'图 {idx + 2}  {filename.replace(".png", "")}（路径：{route}）\n{caption}', max_width_cm=15.2)

        story.append(Paragraph('4. 测试报告', styles['CNH1']))
        story.append(Paragraph('作品运行环境采用 Docker 开发模式，服务由 compose.dev.yml 管理。验证时重点关注容器服务状态、健康检查接口与关键业务页面是否可访问。', styles['CNBody']))
        story.append(make_table_pdf(TEST_ROWS, [3.5, 6.3, 6.2]))
        story.append(Paragraph(f'验证截图均来自真实运行站点，正式访问地址为 {PUBLIC_URL}，本地演示地址为 {LOCAL_URL}。', styles['CNBody']))

        story.append(Paragraph('5. 安装及使用', styles['CNH1']))
        add_pdf_bullets(story, styles, [
            '准备 .env 文件，并根据环境填写必要的配置项。',
            '在项目根目录执行：docker compose --env-file .env -f compose.dev.yml up。',
            '等待 web、postgres、redis、worker 容器启动完成。',
            f'浏览器访问 {LOCAL_URL} 进行本地演示。',
            f'正式展示可使用公网地址 {PUBLIC_URL}。',
            '根据页面导航依次演示首页、题库广场、题库详情、我的题库、数据中心、考试中心、论坛和后台。',
        ])

        story.append(Paragraph('6. 项目总结', styles['CNH1']))
        story.append(Paragraph('Sak-AI答题助手围绕“题库资源管理 + 学习闭环反馈 + 社区互动 + 运营支撑”构建完整 Web 应用，不仅实现了刷题功能，还将题库组织、考试训练、数据可视化、论坛交流和后台管理整合到同一系统中。工程上，作品采用单仓全栈架构，保证 Web 与小程序共享语义，降低维护成本；部署上，Docker 化运行方式提升了演示和复现效率；应用上，AI 题目解析能力让系统更适合面向教学与学习场景的持续扩展。', styles['CNBody']))

        story.append(Paragraph('参考文献', styles['CNH1']))
        add_pdf_bullets(story, styles, [
            'Flask Official Documentation. https://flask.palletsprojects.com/',
            'SQLAlchemy Documentation. https://docs.sqlalchemy.org/',
            'Docker Docs. https://docs.docker.com/',
            '微信开放文档 - 小程序。https://developers.weixin.qq.com/miniprogram/dev/framework/',
        ])
        return story

    build_pdf(DESIGN_PDF, '软件应用与开发类作品设计和开发文档', [PROJECT_NAME, DATE_TEXT], body)


def build_runtime_docx() -> None:
    doc = Document()
    configure_doc(doc)
    add_paragraph(doc, f'{PROJECT_NAME}-运行网址与答辩说明', style='Title', align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, size=18)
    add_paragraph(doc, f'整理日期：{DATE_TEXT}', align=WD_ALIGN_PARAGRAPH.CENTER, size=11)
    doc.add_heading('一、作品定位', level=1)
    add_paragraph(doc, f'{PROJECT_NAME} 是一个面向备考学习与题库运营场景的 Web 应用。系统以 Web 为主应用端，覆盖题库广场、个人题库、刷题、考试、数据中心、论坛与后台管理等核心流程，并通过微信小程序复用同一后端与数据语义。')
    doc.add_heading('二、运行网址', level=1)
    add_bullets(doc, [f'正式访问地址：{PUBLIC_URL}', f'本地部署演示地址：{LOCAL_URL}', '正式答辩建议优先展示公网地址，本地地址可作为现场网络受限时的备选说明。'])
    doc.add_heading('三、建议答辩演示顺序', level=1)
    add_numbered(doc, [
        '首页：说明系统定位、主要入口与学习工作台。',
        '题库广场：展示公共题库与用户公开题库。',
        '题库名片详情：说明题库简介、参与人数和加入方式。',
        '我的题库：展示个人题库整理与资源沉淀。',
        '数据中心：说明学习可视化与复盘闭环。',
        '考试中心：展示模拟考试的入口与使用场景。',
        '论坛：说明学习交流与社区互动。',
        '后台仪表盘：展示管理与运营支撑能力。',
    ])
    doc.add_heading('四、主要亮点', level=1)
    add_bullets(doc, [
        'Flask 单仓全栈架构，工程边界清晰。',
        'Web 与微信小程序共享同一后端语义。',
        '公共题库、个人题库、刷题、考试、数据中心、论坛、后台形成完整闭环。',
        'Docker 化部署便于演示、复现与交付。',
        'AI 题目解析能力增强学习辅助体验。',
    ])
    doc.add_heading('五、截图说明', level=1)
    for idx, (filename, route, caption) in enumerate(SCREENSHOTS, start=1):
        add_paragraph(doc, f'{idx}. {filename.replace(".png", "")}（{route}）：{caption}', size=11)
    doc.save(str(RUNTIME_DOCX))


def build_runtime_pdf() -> None:
    def body(styles: StyleSheet1):
        story = []
        story.append(Paragraph('一、作品定位', styles['CNH1']))
        story.append(Paragraph(f'{PROJECT_NAME} 是一个面向备考学习与题库运营场景的 Web 应用。系统以 Web 为主应用端，覆盖题库广场、个人题库、刷题、考试、数据中心、论坛与后台管理等核心流程，并通过微信小程序复用同一后端与数据语义。', styles['CNBody']))
        story.append(Paragraph('二、运行网址', styles['CNH1']))
        add_pdf_bullets(story, styles, [f'正式访问地址：{PUBLIC_URL}', f'本地部署演示地址：{LOCAL_URL}', '正式答辩建议优先展示公网地址，本地地址可作为现场网络受限时的备选说明。'])
        story.append(Paragraph('三、建议答辩演示顺序', styles['CNH1']))
        add_pdf_bullets(story, styles, [
            '首页：说明系统定位、主要入口与学习工作台。',
            '题库广场：展示公共题库与用户公开题库。',
            '题库名片详情：说明题库简介、参与人数和加入方式。',
            '我的题库：展示个人题库整理与资源沉淀。',
            '数据中心：说明学习可视化与复盘闭环。',
            '考试中心：展示模拟考试的入口与使用场景。',
            '论坛：说明学习交流与社区互动。',
            '后台仪表盘：展示管理与运营支撑能力。',
        ])
        story.append(Paragraph('四、主要亮点', styles['CNH1']))
        add_pdf_bullets(story, styles, [
            'Flask 单仓全栈架构，工程边界清晰。',
            'Web 与微信小程序共享同一后端语义。',
            '公共题库、个人题库、刷题、考试、数据中心、论坛、后台形成完整闭环。',
            'Docker 化部署便于演示、复现与交付。',
            'AI 题目解析能力增强学习辅助体验。',
        ])
        story.append(Paragraph('五、截图说明', styles['CNH1']))
        for idx, (filename, route, caption) in enumerate(SCREENSHOTS, start=1):
            story.append(Paragraph(f'{idx}. {filename.replace(".png", "")}（{route}）：{caption}', styles['CNBodyLeft']))
        return story

    build_pdf(RUNTIME_PDF, f'{PROJECT_NAME}-运行网址与答辩说明', [DATE_TEXT], body)


def build_sections_docx(path: Path, title: str, sections: List[SlideSection], preface: list[str] | None = None) -> None:
    doc = Document()
    configure_doc(doc)
    add_paragraph(doc, title, style='Title', align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, size=18)
    add_paragraph(doc, f'作品名称：{PROJECT_NAME}\n整理日期：{DATE_TEXT}', align=WD_ALIGN_PARAGRAPH.CENTER, size=11)
    if preface:
        doc.add_heading('说明', level=1)
        add_bullets(doc, preface)
    for idx, section in enumerate(sections, start=1):
        doc.add_heading(f'{idx}. {section.title}', level=1)
        add_bullets(doc, section.bullets)
    doc.save(str(path))


def build_sections_pdf(path: Path, title: str, sections: List[SlideSection], preface: list[str] | None = None) -> None:
    def body(styles: StyleSheet1):
        story = []
        if preface:
            story.append(Paragraph('说明', styles['CNH1']))
            add_pdf_bullets(story, styles, preface)
        for idx, section in enumerate(sections, start=1):
            story.append(Paragraph(f'{idx}. {section.title}', styles['CNH1']))
            add_pdf_bullets(story, styles, section.bullets)
        return story

    build_pdf(path, title, [f'作品名称：{PROJECT_NAME}', f'整理日期：{DATE_TEXT}'], body)


def build_support_docs() -> None:
    build_sections_docx(
        PPT_DOCX,
        f'{PROJECT_NAME}-答辩PPT提纲',
        PPT_SECTIONS,
        preface=['建议总页数控制在 8-12 页之间。', '讲解时可结合截图与实际网站进行同步演示。'],
    )
    build_sections_pdf(
        PPT_PDF,
        f'{PROJECT_NAME}-答辩PPT提纲',
        PPT_SECTIONS,
        preface=['建议总页数控制在 8-12 页之间。', '讲解时可结合截图与实际网站进行同步演示。'],
    )
    build_sections_docx(
        SPEECH_DOCX,
        f'{PROJECT_NAME}-答辩讲稿要点',
        SPEECH_SECTIONS,
        preface=['建议全程控制在 6-8 分钟。', '每一页都尽量遵循“问题-方案-价值”的表达顺序。'],
    )
    build_sections_pdf(
        SPEECH_PDF,
        f'{PROJECT_NAME}-答辩讲稿要点',
        SPEECH_SECTIONS,
        preface=['建议全程控制在 6-8 分钟。', '每一页都尽量遵循“问题-方案-价值”的表达顺序。'],
    )
    build_sections_docx(
        VIDEO_DOCX,
        f'{PROJECT_NAME}-演示视频脚本',
        VIDEO_SECTIONS,
        preface=['建议视频总时长控制在 6-8 分钟，最长不超过 10 分钟。', '录屏时优先使用正式地址，如网络受限则切换本地部署地址。'],
    )
    build_sections_pdf(
        VIDEO_PDF,
        f'{PROJECT_NAME}-演示视频脚本',
        VIDEO_SECTIONS,
        preface=['建议视频总时长控制在 6-8 分钟，最长不超过 10 分钟。', '录屏时优先使用正式地址，如网络受限则切换本地部署地址。'],
    )


def build_text_files() -> None:
    MATERIALS_TXT.write_text(
        f'作品名称：{PROJECT_NAME}\n'
        '1. 网站截图：来自本地运行站点 http://127.0.0.1:8000 的真实页面截图，保存在“2 创作素材/网站截图”。\n'
        '2. 系统架构图与核心流程图：由本次整理脚本基于作品实际架构与业务流程自制生成，保存在“2 创作素材/插图”。\n'
        '3. 答辩演示PPT：由本次整理脚本基于作品截图、架构图和功能亮点自动生成，提供 pptx 与 pdf 两个版本，保存在“3 作品”。\n'
        '4. 作品运行演示视频：由真实页面截图自动生成 mp4 演示视频，覆盖首页、题库、考试、论坛与后台等核心页面，保存在“3 作品”。\n'
        '5. 源码包：来自当前仓库源码，保留 Web、后端、小程序、Docker 与文档目录，排除缓存、运行数据、备份和版本控制元数据，保存在“3 作品”。\n'
        '6. 正式文档：依据竞赛模板与附件要求整理生成，统一保存在“1 文档”。\n',
        encoding='utf-8',
    )
    WORK_URL_TXT.write_text(
        f'作品名称：{PROJECT_NAME}\n'
        f'正式访问地址：{PUBLIC_URL}\n'
        f'本地演示地址：{LOCAL_URL}\n'
        '说明：评审优先使用正式访问地址；如现场网络受限，可说明使用本地部署演示地址。\n',
        encoding='utf-8',
    )
    VIDEO_README.write_text(
        f'本目录对应 {SUBMISSION_ID}-04作品演示视频。\n'
        f'已包含正式提交用《{PROJECT_NAME}-作品运行演示视频.mp4》。\n'
        f'同时保留《{PROJECT_NAME}-演示视频脚本》作为补充说明材料，便于后续重录或替换。\n'
        '如学校下发正式作品编号，请同步替换总目录名与四个子目录名前缀。\n',
        encoding='utf-8',
    )


def should_exclude(path: Path) -> bool:
    parts = set(path.parts)
    excluded_names = {'.git', '.venv', '__pycache__', '.pytest_cache', 'var', 'backups', 'tmp', '.DS_Store'}
    if parts & excluded_names:
        return True
    if 'output' in parts:
        return True
    if path.name.endswith(('.pyc', '.pyo', '.log', '.sqlite3', '.db')):
        return True
    return False


def build_source_zip() -> None:
    include_roots = [
        ROOT / 'app', ROOT / 'static', ROOT / 'templates', ROOT / 'miniprogram-1', ROOT / 'docker', ROOT / 'migrations', ROOT / 'docs', ROOT / 'scripts'
    ]
    include_files = [
        ROOT / 'README.md', ROOT / 'README.zh-CN.md', ROOT / 'README.en.md', ROOT / 'requirements.txt', ROOT / 'compose.dev.yml', ROOT / 'compose.prod.yml', ROOT / 'run.py', ROOT / 'AGENTS.md'
    ]
    if SOURCE_ZIP.exists():
        SOURCE_ZIP.unlink()
    with zipfile.ZipFile(SOURCE_ZIP, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=8) as zf:
        for base in include_roots:
            if not base.exists():
                continue
            for file in base.rglob('*'):
                if file.is_dir() or should_exclude(file):
                    continue
                arcname = file.relative_to(ROOT)
                zf.write(file, arcname.as_posix())
        for file in include_files:
            if file.exists() and not should_exclude(file):
                zf.write(file, file.relative_to(ROOT).as_posix())


def verify_outputs() -> None:
    required = [
        SUMMARY_DOCX, SUMMARY_PDF, DESIGN_DOCX, DESIGN_PDF,
        RUNTIME_DOCX, RUNTIME_PDF, PPT_DOCX, PPT_PDF, PPTX_FILE, PPTX_PREVIEW_PDF,
        SPEECH_DOCX, SPEECH_PDF, VIDEO_DOCX, VIDEO_PDF,
        DEMO_VIDEO_MP4, SOURCE_ZIP, WORK_URL_TXT, MATERIALS_TXT, VIDEO_README, ARCH_IMG, FLOW_IMG, PACKAGE_ZIP,
    ] + [SHOT_DIR / item[0] for item in SCREENSHOTS]
    missing = [str(p) for p in required if not p.exists()]
    if missing:
        raise FileNotFoundError('以下文件未生成成功：\n' + '\n'.join(missing))
    for pdf in [SUMMARY_PDF, DESIGN_PDF, RUNTIME_PDF, PPT_PDF, PPTX_PREVIEW_PDF, SPEECH_PDF, VIDEO_PDF]:
        pages = len(PdfReader(str(pdf)).pages)
        if pages <= 0:
            raise RuntimeError(f'PDF 页数异常：{pdf}')
    if DEMO_VIDEO_MP4.stat().st_size <= 0:
        raise RuntimeError(f'视频文件异常：{DEMO_VIDEO_MP4}')


def main() -> None:
    ensure_dirs()
    build_diagrams()
    copy_submission_assets()
    slide_images = build_ppt_slide_images()
    build_summary_docx()
    build_summary_pdf()
    build_design_docx()
    build_design_pdf()
    build_runtime_docx()
    build_runtime_pdf()
    build_answer_ppt(slide_images)
    build_answer_ppt_pdf(slide_images)
    build_support_docs()
    build_demo_video()
    build_text_files()
    build_source_zip()
    build_package_zip()
    verify_outputs()
    if TMP_DIR.exists():
        shutil.rmtree(TMP_DIR)
    print(f'已生成交付目录：{OUTPUT_ROOT}')


if __name__ == '__main__':
    main()
